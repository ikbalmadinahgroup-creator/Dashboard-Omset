"""Dashboard Omset MFlash

Dashboard Streamlit untuk memantau Omset, Iklan (Meta Ads), Walk-in,
6 Pilar MFlash, dan Kontribusi Marketing Corporate vs Sales Retail
di 18 cabang MFlash. Termasuk styling tabel Walk-in (kotak + warna),
export tabel Walk-in & Scoreboard ke JPG/PDF, insight otomatis, dan
export laporan lengkap ke PPTX/PDF. Scoreboard mengikuti persis format
& rumus pada sheet "Scoreboard" di file Excel master (target kuartalan,
EXPECTED VALUE berdasar hari berjalan dalam kuartal, % PENCAPAIAN =
S/D HARI INI dibagi EXPECTED VALUE). Tabel Walk-in per cabang bersifat
KUMULATIF dari awal kuartal (1 Juli/Okt/Jan/Apr) sampai Tanggal Acuan
yang dipilih (konsisten dengan S/D HARI INI di Scoreboard). Loader data
Omset & Walk-in memakai pandas.read_excel (bukan openpyxl read_only)
supaya tahan terhadap file export MFlash dengan metadata dimensi sheet
yang tidak akurat. Loader Iklan mem-buang kolom duplikat sebelum
digabung (pd.concat) untuk mencegah pandas.errors.InvalidIndexError.
"""

import base64
import calendar
import io
import os
import re
from datetime import datetime, date, timedelta

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

import openpyxl
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import requests
import streamlit as st

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib import colors as rl_colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

LOGO_BASE64 = "iVBORw0KGgoAAAANSUhEUgAAASwAAADUCAYAAAAmyx61AAAuOElEQVR4nO3de3xV1Zk38N/zrL3PyUlCIDcuigIB0SLiJQlQGYvW1mI7fWvbwUoSsLaOTm2tCt6mtkOZttrqCFqr0zq1rUJAzbS+tdV2pq2XvtYCId6LlqsoipAb5HZyztlrPe8fJ8EASUhCknOQ5/v5xA+es7PXs5N9nqzbXgtQSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUOtpRqgNQKTR3rjd2X1GuNW05jjzPM16s3U/s2/fc6sZUh9ZJAGq655S8UMKOJEr44owNnDRlt5zQQMueDVIdnxpemrCOMeOKFxQkyHwMkE8I4UwAx0GQTQQjkASAJgLvEGAdkzwZ5vj/27m2KjqcMcpPxmXG2iJzAXzSAbNEcCIIIwTiQcgxSQtA7xLTiyzy+1jg/THn+k11wxmjSg1NWMeIguLyccT0NQEWEfN4gCDiAAggXQ4kgEAAcfJ9kVcFdF/EtD841Ilr1x1jskb52ZcB8hXDNM03QOAA6wTSJUYiwBDBGEAESASyUwQPWQl+lL34rV1DGaNKLU1Yx4D8krKFxPw9InOCiMUBn/7DIQYRQ5xdT84uqa1Z89xQxNh8x+Rz/ZDcGfborMACcdv3GD0mhDwgHuAta+03Mxe/uXIoYlSppwnrA2zKlHnhvbn5dxLxVwUAxA34XMQGIq5NnNxQv6HyvkELEkDz8qJrQoa+bxgZsaAfyfQgviEQAYlA7t3JvGTqNVtigximSgOasD6gJsy9NKOl1T7IxlwsbpD6polAIFgX3NKwYc2tg3HK5jsnfTsS4qWBE9iB59P9iIBIiNAWd4/saXRfnLRsR/uRn1WlC051AGpIUEtr8MNBTVYAIAIRgWHvewUlZVcc6ema75z0tUiIlybs4CQrINnabYsJMkP8hdGj+O7BOatKF1rD+gAqKCm7goz3E3F2aAoggkBaIXxuffXKDQM5xb4Vkz8cZvkTgMhgJauuiICwR4jG3T9nL97+08EvQaWCJqwPmMLiRZPF2HUA5ferc72fiA3EBWuzMv3zdjz7YL+aXXLn+EiU/T+HfS5pTwxdjB4DTlAH4tmRa7ZsHbKC1LDRJuEHyfz5xrG9h8gMabICAHEWxP7slrbg5v5+bwv7/xoJDW2yApJTIsIeFThr75FHYYa0MDUsNGF9gBS+GbqK2Vw4ZE3Bg4izIOCmwpIvnN3X72lbXjQnxHTDUCerTtGEIBLiC1t2Tr5qWApUQ0qbhB8QuSUV0w3hORBGDnXtqitiA2eDlxHKOKf++Z8193as/PTkEW3NiefCHs04kukL/WUYEMG+IIF/GHHDtteGrWA16LSG9QEwZcq8MEPuJeZhTVZAspbFxj+d4rGlhzs22pxYlhka3mQFADbZNBxJBvduuntKeFgLV4NKE9YHwN5RuUvYeB8ZrqbgwcRZgOnreSVlH+vpmJblEy8wjKuHqyl4sGhCkBWmjxxn7ZKUBKAGhTYJj3L5JeWlxPQMgMzhrl0dgBgQ+XvcxM5uWlvV0PWtffecku8Fsb+GPD4pPsy1q65M8m5vs5Bzs67dXp2yQNSAaQ3rKDZmRkUWAfcRcWqTFQCIA7E5OWRDt3V9+eml8LwgfldmhklpsgIAK4DvUaYI7tt1x5islAajBkQT1lHMhtw3yXglqWoKHkxcABBdXjhz4fzO187FXBDwSizu9po0uNvaE4KsEJfkmKxbUh2L6r80uIXUQBSULJwL5sXpkqw6ERkW52YBgCwFv52z1c+8btsdgchCAuKUBp0Q7QmBMVjSdOfEuamORfWPJqyj0MjTLx0F2PsIHDpwMatUIzgbRGNxvl9+Mi6zbVTRr8f64Zda755cnH3t9t8mnFSHTOozlhPAYwp5TPfKT4pGpjoe1XeasI5CfijxHTL+NJH0ql0lh3Ak2vz2cbtBnAfgH/1sMxVOSjqO2MOpz1cAgFggiIT41NZWfDfVsai+S5PbR/VVwayFn4TD44CY9KpddSInkAsaqiv/1Lpi8iJiN6XdueWj9ma1tY1sfynk0Yf6szjfUEquXAqbCPB/spdsezLV8ajD04R1FBl7xmWFCT++loiLjmQxviFFDHHyUsD43L71q7Z3viwCar1r0jcyQ/zdaDw9EhYAhAwhbmWbZZ6dc82W2lTHo3qnTcKjSMKL3cHspW+yAgBxYKYzPCfP55eWPVAwc9EPC4rLFxBBsvZG7miLybMRP33+TsatIDNERcba21Mdizq89LlzVK8KSssvJuKHRdzR8TsjAhEna1w2sVec+3B9zZo3mpdPPtU38hyAUUOxDtZAEAEek8QCe8mIxW8+mup4VM+0hnUUyJtZPh7ACjma/sCIQJyF2ASIzSgi+hGKr/BHLN76t0Qg3/LTYLSwU8ecW/KYl7fePWV8isNRvdCElf6IRO4i9o5L66ZgL8RZkPHPz6eWawAg+4Tt/9ked09GQumTtBJWkOHT8c66u0SOoj8MxxhNWGmuoKT8Mmbv84O6NnsKiFiAeGlBySVn0sWwDHt1LO5qvTS6A6NxQcSnz7feNemyVMeiupdGt4s6WOHsiikg+oEcpTWrA4iAmbMBvm/87PmRyOK3tlmhGwxT2lRnBMlNWw3RD/beOXlKquNRh9KEla7mzvVc4O4h5oKUP9g8SDqahrOj1r8ZALIWb3uoPXCPplPTsHNZZZ/cPU8vhZfqeNSBNGGlqYLW477KxpuXbs8KHqnkssp8Y2Fx2RwChJy3uD3udqZTJ3zHssrzSkYV6bLKaSZ97hK13+jistMc83Mg5HxQalddJZdVti8jFD6n/vmfNbfeNenzHnOVdULpcrnJZZWliUFzwtfqssrpQmtYaWbKlKvDlnAvMX8gkxXQuayydzri0WUAkHXt9l8Ggfw8nSaUWgeEPM5JCO4TXVY5bWjCSjN7cxuuZ+Of80FrCh5MnAWRuTqvuOLjAGCdd1M04bakw2oOnZJrZ9E5rc5en+pYVFL63B0KecULZrIxT0OQ2fcHmwnp+RB0HxBDnNsUJz67uXplfcudk+f5vvzGOnjpUrlkAojQZgM5L2vJ9vWpjudYpzWsNDFmRkUWMyeXO+4uARGD2IDYA1HHnqACC7hEcn1iev99NsnnTdKdOLDxpobE3QYA2Uu2/j4WyH2D2TQkev9rf7E9vN9dqU4AnykTpMsqpwMdtk0TNuS+RewXHzBBlBhEDHGBg8hmca5GgJcgtDkO2tVs/WZYsiD4GV4sN0LBCUyYBlAxgDOYeSyIIOKQrv1h4gIQ05fzSst/11Bd+dgIY77VFrMfDfs8faDbgXUmJ2uBRECIJwiBJVhLEEkmLOo4jlngGcD3BL4v8IyA6MAfV3sgyAxzsYtlfgtAv3e6VoPnKPgz/MGXX1x+Lhn+H4gkVxAlBhFBnNsGkSon9FhDzTkvAFcmOr9HBIzf5WWjMR6Czw6lTW3eJLS/3/O1NGdEyZZzQ4QFBPoksckRsemZuIgBcW/BYXZdTeWuphUTPxJi/oMIQq4f4RIBzgHtMUY0RognGLa7rsAeWtHMgOcJImGHSFjgecnsJkg2DZkQjzv6xIjrtj4zsAtVR0oTVoqNPP3SUb6f+AuxmSYiIDYQF7wO8Io6co9gfWUTAMiPJowNLP9DAD5HRGaIyHgCckQoBECEEGWgjli2ElAdzpCnMq56c13MAph6+aT8kdGrCfRlYs5Jxw59Yg/OJtbUb1hdDkBaVhTdmxWiq9r6sHZWZ6JqjTJao4wgoP2v91dnPmcGIhkO2ZkOvicQAcIeIRbIxjjsnNzrduzt/9nVkdKElWIFxWW3k+ffABE4cS0kcoeV+N2NNVX7ACB6V9H5BPmygD4eMlTAnBxydwK4LrUlAiVrAZysDbQnBET4GwSPhLP3PUCX178bmn7pKTmR4DYic5Eg/ZqJRAwr7nMN1ZWPtd1VdCIELzEjt7dlaIiAaDujqYWRCJIdUYN1U4skf57ZmQ7ZmTaZxHxCS9zdPuK67TcNUjGqHzRhpVDh7IopzsqLbPxsZxMvinP/0lCzZj0A7PuPyWdnePItIczrWBUT/Vk/igB4huAbIBZInQh+kjFx2w/oIjTnl5b9C8H8gAg56fScIpGBc8HL2Vn+7B3PPtjesrzovqwwfaW7WlZnrWpfi0FblCEydOMMIkDIF4zKsYiEBE7QDEtnRZZs3TI0Jaqe6ChhCom1VxovI1ts4pcJE/9YQ82a9XLn+Ej07qLbwp485Xs0z0nyUZH+LnYnSC6Z0vFhL8jw6Zb4m0XroysmfqK+evWPrciFIvImmc5RxdTfCiIWzN7pLdH4PABgopWxhNiDExEREASEukYPrW28/7WhQgTEE8nyWqKMjBCNCMhdMXQlqp6k/i49RhVOm58tQhUuaH+orjVe1rS2qqHtrqITYyb0RIZHNwsQjiZkUFpt1gFtcQEbOoWZf9N2V9HNezesel6cvRDOVTnrHoW4NygNklayQ5wqACDC9ELgsMXrstUOUXLkr26vQTxBwzZ7o3PksGGfwb4WBhM+L3fM0GkOwywN7tBjk42EPwzgz3Uu63JsrIo3/ceEUwj4n7BP57XFBf0ZHeureCBwAj/i020tK4puq69Z80bt+pUX129Y9YW6XaPPck6WpzppJadgYM7I08py6ZotMRFZ27lmVmfNqn6vQRAMX7I6WF0jIxY3k6LhljNSE8GxSxNWqnhoipO5CjX3JxpXTJjoe+bXYY9O6cuo2JFwgo65SNIEAK0rJs4KfjT5a/KtX4fqJ8VuFGc3pDRpiYCIxnCGnAIAzPRC58ROjwh7m1KbrIBkc7s96lFTM5cc9mA1qDRhpUjDulXrmqtX1suPpmWHwJUZPk2NJoZ+1I4IaEuIE6LfyNJpIRF+1GSZe9pacROqqixAz6a6P4vYEAtPBgCIbLcu2VJsjmJFNE57TRrctc4RmprNSamO41iTBr/6Y1trLHprZpjPHuqa1X4CGAIZh5HAxgAkNbbN7nWQv3UcUDB4zyYe+LhQ8rGivtxyBJCM6ThDgzEE61CVf/PWxQT5fjgN1lUWEQiSMarho4/mpFDz8qKPegZXtQ/jxqICIOwRtSfkKlqGv8ijkUva3k4Ujliy9R1Mmx8SorMGpbVFDDgXOGf/AJJnyKIeRMcL4QIimpMMpufrJqHkki4Mao+7nSTmOgCIh70ViAfzIh6f2z7AR3cGgwAgQihlARyjNGGliCydFmpB+62+IdM+DE3BrtoTgrBHZS3LJ+2hizdeB+AdAMDGqjiVlH1PRNbgCOboETEcZBMEV9ZvqHzmoLf/vaC0/GIQ7gNxfo87AYlEAcCK9QDv+hFLNr8DANOXbYxv+tcpX41b+ovh1O1tSAAYFE1N6ccuTVgp0pYb/XTE8KzhTladiABibD349boNqx/JLymfx8b74oB26iGGiLxLSHy6ruaRTXnF5RewoSshmAKgQUSq6qor78svLqsllidAFDm0piUAyS4AyA5GVuOGV9q6vjv1ti0bN//r1G9mGP6RHYrh1D4gIgiSMarhk/rOgGOUEzohVZ3HmSFCW9w9dnvj9vu6e98PQjc6F2wdSOd78qFt+fe66kc25ZeUX8WGnyQynwPRDGI+lz3v3vySsqr6mtXPEOSm7vq0xFkH4U0AQDe80krddKpN2XL6j9sT8kSGn8pb2G1MYeHHJE1YKcJGnm9PDP++874htCdkl3F8zbJl6LZB9d5LP68lkauT6231AxGcDeqDiDw6uqSsCITbARhxASCuYyfoAOz5/1RQUvbPtdWV94gNKom7VPSJAchOCYc391pUVZV1cF9PBG5P14mlw4EIaE+4wBGqh7VgpQkrVeoTiVetw2ZvGJcE7iwpsFicef3Wt3s7tm7D6t+J2PsOSCaHPT8BhB37nlvdKESfYvayuu2jSjYBLwaARAauFhe8QpxclJCIAcHT9c//rPlw5Z38/c3bnHM38iA+8NwXHhMc5PWE72sNa5hpwkqRE5fsjBLJw74ZvjIjIULMysoRS7Y93JfjnSS+JTb42/4VTvvBAfk9vikCgEYBwL7nVjfC4TIRty/Z/+VEiFb2tZzJ39/8UCyQR4ezaegzgUTWTF+2MT5shSoAmrBSiiw/EE1Iw3D0ZYUMIRqX7dlh9HlDheQSN+6rAhfvSx1GICDBhLyZ5TkCvNzjfK5ks++Nzv+tq6l8gcRdS2wg4p6vz9z5bF9jJEDCYVzXnnA7h6O2apjQFne1CTa/GPLC1CE0YaVQ5vVb33bWrQh7Q/tBIwIcxDmSr9NV2/b053vrNqx5VkTu7Gyy9UoEZLx8EilrcFm/dTZ4+ZAmJTHE2Zglubfry7XVq38hQfweBn0Xzz7br+HJCcv+/q5zuBYCGepHdjp29bl92q1v6AhhCuh6WCkmPxmXGW2L/CnDp9lD9WhOZpjQ3C4/ylm87eqBfP+44isyE9zyFLE367CrlRJBgPc84Gwr5INcJZEp6bzVxLk9Im5J/YbKVQOJpTebbp56f1bY/HNbfGgmZ0V8RjRh/+yFYp+YtGxH+5AUonqlNawUoyt3tQXWXhYP5N2hqGklpzDIMyNs9oA3T9hVc3+bOPmiOPfuYR+tEQETjw2ce7DORbZncPwjztrPOJEb4OxlVuIlQ5GsAABhd300bp+PDEF/VnJ5ZPeWsPuSJqvU0RpWmmhbMfnDhuVXnqGxgzWZNDNMaI/LWhu3n82+ccd7R3q+wpKKsx3hV8w85nA1reQa7fEf1m9Yc82Rltsff//GScd7ML/O8Ki4LTE4Na0Mj5Bw8o4TXDTl1r9vGJSTqgHRGlaayLxu618Tzl4YD+TVzHByffaB8jhZs2qPy68TgffpwUhWAFC7YdXz5OSTIvJasm+q5yDFBWD2rs4vKa8YjLL76uRbN78D8T7VHrgnIz7jSAY0mIDMECNu5UXnaJ4mq9TThJVGsq/b8VLchM6LxeU/CUhEQv1LXJ2JSgT10RhuyGjc9k8512+qG8wY62oqX/Di7R8VF9wPQqK3xCUQIqIfjp51yYzBjOFwJt/22u6obz4bS9h/BdCY6TNMP36QTMn+KiaKtSfknlYXOn/KbW+8NnQRq77SJmGaavuPyWcbX651ggszfMoWSS51bA/eKYeTicoJEA/kPZA8Aph7ItdsOeQ5wcFWWFw2R5ivFWAes8kGBCKCrtMZiA0kCP5GFP9YbXXVoNT0+mPjzadMDbF8nQn/5DONISJYJ7CCjliTjxMxAYaS/44FthmgJwLn7jr5+5vXDXfMqmeasNJc+11FJwnJBRCaC+AUERQKEAHgCNJKoF0CvALgT5bkqRHXbt893DEWzFpwklj+OJjmksiHRFBIhAwkM1grsVdrXfDdhg2rfzXcsXXavnTaWBfY8yFyPgSnOdA4EWSBQASJAqglotdJ5BnL7n+nfm/oE77qP01YRxF5eq6HF7bmIJSV0ZqISjzuteXetK0pOZMgTcxd6uXENub4MZNBbMQi3tZYU9WEwVsV8IiJgLbdXJRjTCiTfCGCaT8BhU20rH/zv5RSSimllFJKKTVstA9L9aqgeMEnhb2ZEOuB2EKkpn5D5W+QRn1SA5VfekkJ4P8jxHoACQAwG7E2/v8aah7+Q6rjU4fShDVMli4Ff/u4cRldX6Mrd7V1d+y44isy+3LOILuRXWJEl99hHdgPDyiR1D5b1XLwawXF5beQZ757yG3iErfWVq++ZSDlpIsxMyqygrB7iU14ygFrdhHB2USLDWT63hdX70hdhKo7mrCGQWz55FPFyH9BMDZwAAjwCbCQFyMhezld9VYj0PGQMbXcBzbniLOH/d0QoZslFKifCUsIYAGwzTF9tWHdytcBYOQ/lOV67XiD2Iw+8APNgLh9BO+U2uoHh31e1WAZM6NidBCS14kp75A15QVCJGfVVq9+KSXBqR7pJhTDICC3KDPTfDiICkyXFBMK8aS2NnkQwOMAEDOtMwz7l4oTcA/LuQxVO4yNN0GC2HcA/BMAUCtyYJB5SIkiEJGIoG0kgKM2YRH7AupxWQcH4qO+yftBpAlrGLDj38ejbqG1kte5yYtnSBJtsomsebHzuAzhzTEb/IVAxa6H5d4J5IHZ9Lg91gAlH2amcfvLYSME29OHVkg/0CoFNGENg8iSrU/vuqPojNwMjGJLEguEshhSF/ffK7zp7/vXLn+3emX9+NnzP94u3nhyjsCHJgVyyBLnlhLzRTLISUupdKcJa5iMu2HbHgCHXe1z59qqKIBed4zJK15wG4Mu6vVERCD0fQo8sQE5q4+jdCB2+tcgDWnCGk7FV/hHeooxibaQZfk8gXvZhYsg4vYIsLNvZ2XYIL7JSTCoI39jZlRkuUyMEyfjxCGXQb4D4iyugQ2/u9tm7kTN/YmBnj+3+IqRoLYTmGSMADkQJkPSZp1tEArvamxr3Y2NVf3fKCK5dXXW++XMH2n80GRYOc6JCQFoYsJbdZlvbevvcs7qyOgo4TAoKF04FST3QeS4Q3c57hcRogwCFfXW/U7swbnEXfXVqxf3vZv+wNHF3OKFJxqyr4F5xIExE0RcjMjOqKt+ZNPBZ8mbveh4tvZCAPMAnCmCccQUSW4BRp2d9hBIC4G2QfAHFvfgnprVr/YlyglzL81obU18BuAvgKQEwFgQ+9SxmLuIAOIAkSYQ7xTQX4mDH9StW3NArXXsGZcVJvz4RiIqOOR3QgyIfd6B7mfBbBBdCMgJRMwgSsYvLkqgjQR50HP7HthV89tup6iowaU1rOEg9kvsRc4X26fNZ/pwvj4kIWKH5J7vg1Dg4Y09Y35hwve/QdaVE3uFnUvNECSZpCAH5E4CZXfsBj3DCf1LQUnZPaP2Nnx7y5bfx3oqo7C4bE5r1K4g45UCgIhLnlPcoT8S5hyApjGbaWLdnJzZ8+c0ra1q6NPFiAPInG2Izk6W03ENneUl448QUTHIFCcw6gtjZl1asXvdg2/2/SemBkIX8BsGRPK0s7F6EbHi3BF9wTkLEZtsD4pLbpnVTVISGdbac+D5JWxC14JQKC5IjjpKR87slnTsBh0AkCzyQjc35ub9vKdmc2FpxTxh/h0Rl4qzHec/cO2tA0/fcX6bALE5JeRCH+7XBXXuVL3/Og6NXzriJzZzAhf8Kmf2/Lx+laH6TWtYw6C2es3/5JZ84UzfhHOBwdt7U8QjgS0BeDlAOal8WsaX7GfjQetLzHyG9DcOEYhNgNlfkE8tL9UDt3d9O//MS44TyAMgHtHtWvLEnbtOv7+AYNcqV/LfTf2+qL6G7wKw8c8MBfgG0Pd9H1X/acIaJo0bHnkbQK/bww/QywUlCz5Nxv9M1w8zQYa1M3hXzf1t+aUV/wGiVRAkk0hyw9SOWgosiAwxv9/PdBARCwjdUFh66UNdZ9GL4UuZveOStbGuOvut3DaBvAdHAYAcAOOIaEznXorOBQ/Xu6y1A7qwjtHWZN+V67E53hHbZQXF5XfW1VTqnoVDRBPWMBN5vzLQ4/sdelmYb/8xo0vKJlmiUw+sUTgIcHpeafmlBr01DRlWaGfDp4qewrJlRzyMn51pftnalriZTHi62Phb4twTgPuzE2wTUNQzPNK5YBZA/0LEUw5JWiJg9grEBp8E8LP9URLOPfTHRRBIKxNfZtrx5O5XKluTry/lccWb8hJE0wQ4w1na2dAW/y02VvZ7NJLYQGzQLMDfQBQFMIWIT+h2/psIyJg8EXsugDX9LUv1jY4SDpPW5ZM/E/JxQyJwEd+wiwfybKbhW+iaLTEAaFteNMcYfNeK5HTmnpBhG7PyeNbx226ji2Hzi8vPZUPLnCAb6PzQ0Indj3QRDruHIAgdNaD76ybGr0JV1f4q2kBHCfNnlp1H4DNiQg81V6+s767UvJnl41nojyA6+eCkRexBXLCyrrpyUedrBaVl64lMqRz0kLI42U2CMwdSo+l1lBAdycq5xwG6oa565SYAyJk9P88P/G8zm6u7S1rJ0dlgRX115eL+xqP6RmtYw0CWgtvIfdfzzXQCgwjIDFNJW4D/BrAWAITkm6EMc66Nv//hYQaMQ2n725MfBrZuIZalZEIfYWs7ht47R9+6qYiJJJtYh0UA5PK8HaG7G4CNR3qt9etXPw3g6d6OaVhfuTO/pPwBZnP7IR/85DVNPijGvckpEQceR0xjIPhrQWnZ7yH8gjBvZnZvoyn2Xu3GQ1ef6DNiiLOvUWu8vOt5mtZWNaD4iiUF0nouiE87tFkrADBhwOWqw9KENRyWQWQF/hQk5EOBFWOY0B641yXk3uw8RBw9FcTc+YGD3/m59ITgRJ6Petw5Q/5PzgVzIOIf4XyuLgQAsziMGqQTYvzs+ZFoEJpOTDNEUERw+SIU7nyfkv+Z3V1CFQgIyMHcud7+SZkifwHxx3HwRFkRgGgCkX9l56ijCyRKmaH38ksrNhLkKTL0eO3aVVv6Ez8RQ5x7oNukV3N/AiXl/0vMp3XfNHx/wqkafJqwhgEBIsdvX9L61pQHyaewYZJM399CX3tjf5Mpe8m2O2J3Ff3OkM2OBh2/GGaXyeaNrGu2NAFAXfXq744uLvt14HEmOwlDUAk244/sQWgC4ISY9h7RRSKZqNpt6Cvtjr5EJB8iMpycz2kOeUiop4735JswaDmZgGcBAFbMz9kmvsJsDt1xWgTSZXyBiCIgmsREkwD6lDj7bwWl5f9lYvTt3a+sau3LdYhYkOHeNk3t8QkC6rXPUB0pTVjDhC6GBba82Nsx4Wu3HXazzq4zwvNLytp6/HQQo3P2dy9RARA46x5omBj/O9YfrvSejSteUNBu+RFi76PSMf+qp1G9pL7XEBtrVr6VV1JW4UQeYvbGJUfrekrSnZM8O4ukkWT8620oPhnT5l/Sp0d1nHNgbu7pbaIen4lSQ0wTVirNn2/yt4euBPE573ei95HAAOj2UR8ihoj7g0AegfSStYiEhHbWT0r8sWuH+wBQgmkFGf+jYg8djEt2YFsIJNa5eykRZRxyYC8aNqz+45iZFXOs2OsguAhEJxB1rhnWZQPXQxbjS87xIuN/Ni8LZQ3AL/pwObp0TprShJVCedu8YvLMvQTCQAZsk82jbj5byWfhXqpbv/qBPp2out9FHyBvZvmHILhYbLdTv2Li7HIL9xsWbiDTUTWyMh/sfa8/zdnd61dtB/D1kaeVLfXDZoYgOBOEDwGYBOAEiBxH7OVItzPsBSRYgD4lLJWuNGGlELO3Fy7YC/ZHDegEHQ8Td0dAw/e7ta6EPD90cP9Sch6TXVO3ofIbB39L3qyK6oEuAbjv1dWNSHZwPdv52pQp88J1I/LGehx8DODbAcrrmrQ6KnYTp02bH9o4kBUcVFrQhJVCddUrN+UWl33EI3eGhdAho2A9sQAxFRJoMYiOdAWII8aE3O5riASQ1Hb7PdZ9lozX99CXLuX8J7ecR8IF5ElN7dpVW9ElI3U8NL0DwAP5JeVfY+Y86ebk8XizdoofxTRhpVhjshO9T0urHCy/uLyFPf5x3+ZbDR0hqu1uUn4yLl4wunTBY3uq16wFICNKF+aHIV8BcEW3zwV2p/gKv+CJzT8l4y0CAGdta0FpxasgeUlEtjBQb0FCgtFE8jEAMw6eckDJxwa29bYahEp/mrCOZoTxqQ4BAIil2tmgnYgzDug7EgGIxzvhpwtKyl4VoJ3gJhGb4/ucrAAUmtbzQN4iccm+KQJnEdFskJmdHFKQju2D9s/c7y5KgHjVEVymSgOasNKAPArT/HbR5REfJQkLMkaaowm+d9SSrVsAIK+k7HPEPA8i/P56TJIPwif788Hvd1w9T4wg6bJJRt26NZsLSssfJeMtOmSUUBwAChNzCXU8QNzfmJ1IEXMyGXVG1us8roODNT7Exh+vb01U9atglXY0YaWB9ncmnzMiAz8GAM8A8BmBdScB+MeRMysmscgaYi90YIdPTzWJQYtqH1GoCaDsA1feI0DQ4gfZByyGx2xvdBankvGLk6OFB8Uq7y/gR2QgzloQmb7Mx3LCz5G1UWIT6axl9UXnag1ig/9OtOOKrnOwxCUI6G5fx+TloIddiwBABKaXjrCezqkGgS7glwacSEt7QmLorM8kP5O7AcC30i4i3ayUSSD2evwSIHQkMTXWVO0jkVsBJLqeF8ku/++/99LPD+hM37Pu4d0mRp+EDX4GIJo83hz05QFgKy6ohJFZAtlKxu8udg/Z4/ZnpcYNq15z4ipE3AaBtHee69Dzm/fjFAQi7q/OSXld9eQvdIws7sd+ZguAncQHlW98ANJgjdftYAEACGgrIN383A0A6nUDEXVkdMQkTTQtL5qTYTDDipCQNLdH+fG8m7ftA4DckorpzDwHzhpQx6RGEWKCkS6PghBILJEYIQLjxdp1K5870rjyiitmEeMshnhO2LLYl2trVv+lt+/JLb3kVEPeBRCcBcgYJGvyjQJ6zcD+fk/1mr8CQEHJgrlE3gxL71cdjRABsrm2etXvDznx/Pkmf2fmSZRInC7gU4lkogAFBGQk12OmJgDvCPg1hl1bW514Feh5QmzBrAUnkTMfsxAPRAIRMkTOSbCuvvrhXh7NWcr5JZs+xeRN7IzdiBCEWuIJ89i+lx/c2+sPVSmllFJKKaWUUkoppZRSSimllFJKKaWOHYT583XmehrRR3OGjVBBafnlibhfNRwTC8fMqBhtQ7IAhnMQBDV1NVN/D3S/92B+6cIScsHoupo1Tx7uvCNKF+b3tH1XXxTMLD+HErSr9sXkxhC5JRXTjUjIguvISDkBCXKuCR5tgMWUuurKR3s6V96shR8Slzitsfrhg4+hgpKyeWDvTIjdk5Xprdrx7IPtXQ/ILy6/sj1Bj7W+smoPepBbUnGq2W6n1wEPD/R61eDShDVM8mZWzBKHhZ6faAewsvDMiinwaDJYxomj5+uqi7YUlGz+BBlvrBX7h4aW+J78SOhssBxHRPvq1lc+UTiz4nTnhbbV722OFURC0y3iW4n8UkN0PJi3ds5sL5w2PzsIyQ9g8JgAmwk0dfzsjeG2oGyOMd4JCGJ/rq15ZGtu6SWneuyVinMngHnXtGnzQ3sy/c+CSIj53dp1K5/LLamY7nlULDb+1wT5CR+4L1y84B72pCYIvBnGw3g492Igtj7Evrd7/artBbMWnOSEfXZuBkCS4SUe37m2KgoA4nABPDkNwEVYupTNE1vuFpYqcrKVBBlG3CNkrLNBaIKQKxo965IxVvwLRFxLw8T443k7QuOM8PlW7DuANBsx5xWWlkcgtLl2w6rnASC/eMHJIJRZx98wiBf47+TJmFmXTnSQ81xgt9fXVD4DwjTj8x/zSy85sb764Q0FxeXjyFGWsLQS8wXWybuGzKuAe2PMjIrRCV/OMB4dHzhebzgRl7i01r/48Lu5pZecSmSEHM4SQmuDy/otau7v96atqm/0WcJhwk4uIrKXg1CM4it857nPOLj5TniHwN48ZubWCQL24MSxwy15OZFCItwqgj0AlRSUll/snPs84rHjRub5WcK0yDPhiQz6hiN6W5y7omDWgpMAAJnhGYC8Xb+u8vGGdStfF9/bEGuCYUKGgKKOzS0FsxacxOAlDuZNACcRJNid6X8d4JFgbhHn/q2gZOGZDHeNg2kQmJtY7FiIa2Xn7xQxpxO5a5xghwiuY/EmW3FfwdKlDMtXknCmMEWF+bSo9b7Y+XMgoE6A9vzisvMKntj8CSHZA6CJ4AIIxokxZyQCHm2MtEMkIPEyk4kJ5+bvCH+andwoxB4MN5MjC0iRI3pbyH1ldElZEQAEEdoN0D5j7BcJXqQtd292YINvOuEdRLgwr7ji44A0e8aGRfhyALDCM5yRC4TlJgH5MNIcID7JCX3cZcgZTPIV58xbRuz1YnEKefwlFF/hM8yX2SIDzM1GcG6hablouO+tY4kmrGGQN7N8mhDOFvE+w8RTC7ntfACtJPLL+vUPPU2g7VaklERmgEgEkmfgIiC80LBh9R/F2kdEcBoIwmICjiYcxLE4YpA8Vb9+1VMi9AocjQMAy7KXCGM6y6cgWOyyMk4BcAYnn5vLFssfAuGN+nUPPkNCvxZBiIAJcS/233XrVj4JkW0QdwoIo2BjOQBvMvDqRfDOnhcfeqXjzE/Wr1/1lBB2sk87AWoueHLzlwHZxC7IgOOTiEwUkowLAAQIichPQbQAoPMJqGThMIwjIYrDodUYExMRIiJnIdNY5AQh0w7BOCeyimBHs5MLROwoISSvH/S3BHgsAOx7bnVjXXXl1cJcJXCL4mQXgbCzfv0vngLjt0Q4vWMnVun8ABgKDENInKwmcoXscAGDM4kpgDgC47f11Q/+CSS1BvQ3J8jK55YvU3Ij3HEs7kQh026Fxg7PXXVs0oQ1DFjkcwAtM8Y+5IxZLJDzSCgDwD8Xli68WiD5ztEOEMZbsUSgIEhQHIKzCkorvk7E1zjQ7wT0kiP7debwjQBlsNiAkNyUjwSWvOQuEA3rK1+H0Ov5JeV3F5QuXAzQSHLOidCJ1jmBCDPF15NgQkf5C4XQQuA1fuDfUlBa8U0QJlnynwPwjjGhEQSpr5vYvhUACkrKvyQioc6yAUmYmEuI4FER+lqMzC8d0SQiCYtLhIiwf1kXIlh2aATx/Zbsz8WhFdKxNrRITW31yv+pXb/q5eQ1SUwcJoiABUGIxFlD5gQhUydOJjJTmEDJ5leX688tXnhifknZd8i5c0CwsPQsCTIKZi26Bg5ljoLfAQyKJ/YJeE9h6aKbhKnCCQX7z08ygZwLk7UJErIkbAGAgAT53MJMT5Lgyxkm/htApu6P0UHXix9C2oc1DAj+T+uqH3yv4393j5lRcWcQcotIqEo89/cIEj/dWV0VHTOzYo91yPIC/0n2YiMc8wZr3dPi8Ou9L1buAIDc0ks2GaZ4nIL63Lrm1tqxeTsBIBGYn4/KNp0dy1K3ofKuwjMrpjjjchEK/1fd8z9rzi35wnd8plwAf9yzrmr3mBkVN9pMd7LE7C85xzahzZ9CzNXiYAg0tnHDL3aOmVFxi43gFDi8i6oq64rn/5vvmQnWZG6XIEYAkDCZPyRuyyeYM8TZ/5vslJdVo4vLp4vjaMCZ+5dqScS9n42ItMV2rl0dBYAxMyreTGRITSbbRKu1NZ3HRSLeS3v34o19U9uaC9/0TzPEj7vWWK0JZUcSvkwMs/nVuzbSNDLS8jwABDH56aisULQBQGPNyrfzZpb/WIiO89g9vqf64d3Tps1/bU9O1nQieqSxuvK9ccUL7t41xTZiW853CrzW01xgV3tR12hCiCT8rImhQB7bhRH7RgaxrEyvPh7jkR4A+DDL260bw4KTQbRm59qqKObOvaew7fjphvjxIByqG9q76dimy8ukSH7xglOMJ4171j28u7v3C+fOz5aW0NS6msoXhiumvNmLjifrzmWQJ879b11N5a6+fm9+SfnJxDzH2vZfNtZU7RvKOFOtoLj8LGKcFjbxRzsHE5RSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSSimllFJKKaWUUkoppZRSqfH/AXlTRJE7lZLQAAAAAElFTkSuQmCC"


def _page_icon():
    try:
        raw = base64.b64decode(LOGO_BASE64)
        with open("/tmp/_mflash_icon.png", "wb") as f:
            f.write(raw)
        return "/tmp/_mflash_icon.png"
    except Exception:
        return "📊"


st.set_page_config(
    page_title="Dashboard Omset MFlash",
    page_icon=_page_icon(),
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown(
    """
    <style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    </style>
    """,
    unsafe_allow_html=True,
)


# ========================= GitHub Auto-Backup =========================

def _get_secret(key, default=""):
    try:
        return st.secrets.get(key, default)
    except Exception:
        return default


_GH_TOKEN = os.environ.get("GH_TOKEN", "") or _get_secret("GH_TOKEN", "")
_GH_REPO = _get_secret("GH_REPO", "") or os.environ.get("GH_REPO", "")
_GH_BRANCH = _get_secret("GH_BRANCH", "main") or "main"
_GH_ENABLED = bool(_GH_TOKEN and _GH_REPO)


def _gh_config():
    return _GH_TOKEN, _GH_REPO, _GH_BRANCH


def _gh_headers():
    token, _, _ = _gh_config()
    return {"Authorization": f"token {token}", "Accept": "application/vnd.github+json"}


def github_get_file_sha(path: str):
    token, repo, branch = _gh_config()
    url = f"https://api.github.com/repos/{repo}/contents/{path}"
    try:
        r = requests.get(url, headers=_gh_headers(), params={"ref": branch}, timeout=15)
        if r.status_code == 200:
            return r.json().get("sha")
    except Exception:
        pass
    return None


def github_upload_file(path: str, content_bytes: bytes, message: str = "auto-backup"):
    if not _GH_ENABLED:
        return False
    token, repo, branch = _gh_config()
    url = f"https://api.github.com/repos/{repo}/contents/{path}"
    sha = github_get_file_sha(path)
    payload = {
        "message": message,
        "content": base64.b64encode(content_bytes).decode("utf-8"),
        "branch": branch,
    }
    if sha:
        payload["sha"] = sha
    try:
        r = requests.put(url, headers=_gh_headers(), json=payload, timeout=20)
        return r.status_code in (200, 201)
    except Exception:
        return False


def github_delete_file(path: str, message: str = "auto-delete"):
    if not _GH_ENABLED:
        return False
    token, repo, branch = _gh_config()
    sha = github_get_file_sha(path)
    if not sha:
        return False
    url = f"https://api.github.com/repos/{repo}/contents/{path}"
    payload = {"message": message, "sha": sha, "branch": branch}
    try:
        r = requests.delete(url, headers=_gh_headers(), json=payload, timeout=15)
        return r.status_code == 200
    except Exception:
        return False


def github_download_file(path: str, local_path: str):
    token, repo, branch = _gh_config()
    url = f"https://api.github.com/repos/{repo}/contents/{path}"
    try:
        r = requests.get(url, headers=_gh_headers(), params={"ref": branch}, timeout=15)
        if r.status_code == 200:
            content = base64.b64decode(r.json()["content"])
            os.makedirs(os.path.dirname(local_path), exist_ok=True)
            with open(local_path, "wb") as f:
                f.write(content)
            return True
    except Exception:
        pass
    return False


def github_list_dir(path: str):
    token, repo, branch = _gh_config()
    url = f"https://api.github.com/repos/{repo}/contents/{path}"
    try:
        r = requests.get(url, headers=_gh_headers(), params={"ref": branch}, timeout=15)
        if r.status_code == 200:
            return [item["name"] for item in r.json() if item["type"] == "file"]
    except Exception:
        pass
    return []


def sync_data_from_github():
    if not _GH_ENABLED:
        return
    for remote_dir, local_dir in [
        ("data/main", "data/main"), ("data/ads", "data/ads"), ("data/walkin", "data/walkin"),
        ("data/target", "data/target"), ("data/corp", "data/corp"), ("data/log", "data/log"),
    ]:
        for fname in github_list_dir(remote_dir):
            local_path = os.path.join(local_dir, fname)
            if not os.path.exists(local_path):
                github_download_file(f"{remote_dir}/{fname}", local_path)

# ========================= Constants =========================

MAIN_SHEET_NAME = "Faktur Penjualan"
_FAKTUR_SHEET_NAME = "Rincian Faktur Penjualan"

DATA_DIR = "data"
MAIN_DATA_DIR = os.path.join(DATA_DIR, "main")
ADS_DATA_DIR = os.path.join(DATA_DIR, "ads")
WALKIN_DATA_DIR = os.path.join(DATA_DIR, "walkin")
TARGET_DATA_DIR = os.path.join(DATA_DIR, "target")
CORP_DATA_DIR = os.path.join(DATA_DIR, "corp")
LOG_DIR = os.path.join(DATA_DIR, "log")

BRANCH_ORDER = [
    "KLENDER", "CEGER", "BINTARA", "RADJIMAN", "JATIMULYA", "DRAMAGA",
    "CONDET", "JATIBENING", "SAWANGAN", "WARBONG", "CINERE", "CIBINONG",
    "KARAWANG", "JATIWARINGIN", "CIKAMPEK", "CILANGKAP", "PEJATEN", "CIBUBUR",
]
_BRANCH_RANK = {b: i for i, b in enumerate(BRANCH_ORDER)}

BULAN_ID = {
    1: "Januari", 2: "Februari", 3: "Maret", 4: "April", 5: "Mei", 6: "Juni",
    7: "Juli", 8: "Agustus", 9: "September", 10: "Oktober", 11: "November", 12: "Desember",
}
BULAN_MAP = {v.lower(): k for k, v in BULAN_ID.items()}
BULAN_ALIAS = {
    "jan": 1, "feb": 2, "mar": 3, "apr": 4, "mei": 5, "may": 5, "jun": 6,
    "jul": 7, "agu": 8, "aug": 8, "sep": 9, "okt": 10, "oct": 10, "nov": 11, "des": 12, "dec": 12,
}


def order_branches(names):
    return sorted(set(names), key=lambda b: _BRANCH_RANK.get(str(b).upper(), 999))


# ========================= Format helpers =========================

def format_rupiah(v) -> str:
    try:
        v = float(v)
    except (TypeError, ValueError):
        return "Rp 0"
    sign = "-" if v < 0 else ""
    return f"{sign}Rp {abs(v):,.0f}".replace(",", ".")


def format_number(v) -> str:
    try:
        v = float(v)
    except (TypeError, ValueError):
        return "0"
    return f"{v:,.0f}".replace(",", ".")


def format_decimal(v, digits=1) -> str:
    try:
        v = float(v)
    except (TypeError, ValueError):
        return "0"
    s = f"{v:,.{digits}f}"
    a, _, b = s.partition(".")
    a = a.replace(",", ".")
    return f"{a},{b}" if b else a


def format_percent(v, digits=1) -> str:
    try:
        v = float(v) * 100
    except (TypeError, ValueError):
        return "0%"
    return f"{format_decimal(v, digits)}%"


def sanitize_filename(name: str) -> str:
    return re.sub(r"[^A-Za-z0-9_.-]+", "_", str(name)).strip("_")


def branch_from_filename(fname: str):
    """Deteksi nama cabang dari nama file. Cocok persis dulu; jika tidak ketemu,
    coba cocokkan token setelah 'MFLASH' dengan pencocokan awalan (prefix) dua arah
    supaya tahan terhadap nama cabang yang terpotong pada file export (mis. file
    bernama '...mflashklende_...' tetap terdeteksi sebagai KLENDER)."""
    up = str(fname).upper().replace(" ", "")
    for b in BRANCH_ORDER:
        if b in up:
            return b
    m = re.search(r"MFLASH([A-Z]+)", up)
    if m:
        token = m.group(1)
        for b in BRANCH_ORDER:
            if b.startswith(token) or token.startswith(b):
                return b
    return None


def branch_from_sheetname(sheet: str):
    if not sheet:
        return None
    up = str(sheet).upper().replace(" ", "")
    for b in BRANCH_ORDER:
        if b in up:
            return b
    return None


_SERVICE_KEYWORDS = ["SERVICE", "JASA", "SPAREPART"]


def classify_kategori(v) -> str:
    """Klasifikasi Service vs Gadget & Aksesoris berdasar KATEGORI PENJUALAN."""
    if not v:
        return "Gadget & Aksesoris"
    up = str(v).strip().upper()
    for kw in _SERVICE_KEYWORDS:
        if kw in up:
            return "Service"
    return "Gadget & Aksesoris"


def _extract_filename_timestamp(fname: str):
    """Ekstrak timestamp YYMMDDHHMMSS di akhir nama file (format export MFlash)."""
    m = re.search(r"(\d{12})(?:\.\w+)?$", str(fname))
    if not m:
        m = re.search(r"_(\d{6,14})\.\w+$", str(fname))
        if m:
            return m.group(1)
        return None
    return m.group(1)


# ========================= 6 Pilar MFlash =========================

PILAR_ORDER = ["Handphone", "Laptop", "Aksesoris", "Voucher & Perdana", "Service", "Lainnya"]
PILAR_ICONS = {
    "Handphone": "📱", "Laptop": "💻", "Aksesoris": "🎧",
    "Voucher & Perdana": "🎫", "Service": "🔧", "Lainnya": "📦",
}
PILAR_COLORS = {
    "Handphone": "#2563eb", "Laptop": "#7c3aed", "Aksesoris": "#d97706",
    "Voucher & Perdana": "#059669", "Service": "#dc2626", "Lainnya": "#6b7280",
}
_PILAR_SHOW_QTY = {"Handphone", "Laptop", "Aksesoris", "Voucher & Perdana", "Lainnya"}


def _pilar_label(p: str) -> str:
    return p


def _find_pilar_column_index(col_idx: dict):
    for header, idx in col_idx.items():
        if header == "KATEGORI PILAR" or header == "PILAR":
            return idx
    for header, idx in col_idx.items():
        if "PILAR" in header:
            return idx
    return None


def classify_pilar(v) -> str:
    if not v:
        return "Lainnya"
    up = str(v).strip().upper()
    if "HANDPHONE" in up or "HP" == up:
        return "Handphone"
    if "LAPTOP" in up:
        return "Laptop"
    if "AKSESORIS" in up or "ACCESSORIES" in up:
        return "Aksesoris"
    if "VOUCHER" in up or "PERDANA" in up:
        return "Voucher & Perdana"
    if "SERVICE" in up or "JASA" in up or "SPAREPART" in up:
        return "Service"
    return "Lainnya"


def parse_bulan(v):
    if v is None:
        return None
    if isinstance(v, (int, float)):
        iv = int(v)
        return iv if 1 <= iv <= 12 else None
    s = str(v).strip().lower()
    if s in BULAN_MAP:
        return BULAN_MAP[s]
    for alias, num in BULAN_ALIAS.items():
        if s.startswith(alias):
            return num
    return None


def _nan_to_none(v):
    """Konversi NaN/NaT (dari pandas) jadi None, supaya konsisten dengan None dari openpyxl."""
    if v is None:
        return None
    try:
        if isinstance(v, float) and pd.isna(v):
            return None
        if pd.isna(v):
            return None
    except (TypeError, ValueError):
        pass
    return v


def to_date(v):
    v = _nan_to_none(v) if not isinstance(v, (datetime, date)) else v
    if v is None or v == "":
        return None
    if isinstance(v, datetime):
        return v.date()
    if isinstance(v, date):
        return v
    try:
        return pd.to_datetime(v).date()
    except Exception:
        return None


def _to_float_or_none(v):
    v = _nan_to_none(v)
    try:
        if v is None or v == "":
            return None
        return float(v)
    except (TypeError, ValueError):
        return None


def _find_penjual_column_index(col_idx: dict):
    for header, idx in col_idx.items():
        if header == "NAMA PENJUAL":
            return idx
    for header, idx in col_idx.items():
        if "DEFAULT PENJUAL" in header:
            return idx
    return None


# ========================= Marketing Corporate =========================

MARKETING_CORPORATE_NAMES = [
    "ADITYA", "AGUS SETIAWAN", "AHMAD FAUZI", "AKBAR", "ANDRI", "BAYU",
    "DEDE", "DIMAS", "EKO", "FAJAR", "HENDRA", "IRFAN", "JOKO",
    "KURNIAWAN", "MAULANA", "NUGROHO", "PRATAMA", "RAHMAT", "SUSANTO",
    "WAHYU", "YUSUF",
]
_MC_LABEL = "Marketing Corporate"
_RETAIL_LABEL = "Sales Retail"


def classify_penjual_kelompok(nama_penjual) -> str:
    if not nama_penjual:
        return _RETAIL_LABEL
    v = str(nama_penjual).strip().upper()
    if not v:
        return _RETAIL_LABEL
    for known in MARKETING_CORPORATE_NAMES:
        if known in v or v in known:
            return _MC_LABEL
    return _RETAIL_LABEL


# ========================= Loader Data Omset Utama =========================

def _extract_qty_gp(row_dict):
    qty = _to_float_or_none(row_dict.get("QTY")) or 0.0
    gp = _to_float_or_none(row_dict.get("GROSS PROFIT")) or 0.0
    return qty, gp


def _build_col_idx(header_row):
    col_idx = {}
    for i, h in enumerate(header_row):
        h = _nan_to_none(h)
        if h is None:
            continue
        key = str(h).strip().upper()
        if key and key not in col_idx:
            col_idx[key] = i
    return col_idx


def _load_faktur_sheet(path: str, cabang_hint=None) -> pd.DataFrame:
    """Loader untuk format baru: 'Rincian Faktur Penjualan' per cabang.
    Pakai pandas.read_excel (bukan openpyxl read_only) karena sejumlah file
    export MFlash punya metadata dimensi sheet yang tidak akurat sehingga
    openpyxl read_only gagal mendeteksi baris data (sama seperti bug pada
    loader Walk-in)."""
    try:
        sheet_name = None
        xls = pd.ExcelFile(path)
        for name in xls.sheet_names:
            if "RINCIAN FAKTUR" in name.upper() or "FAKTUR PENJUALAN" in name.upper():
                sheet_name = name
                break
        if sheet_name is None:
            sheet_name = xls.sheet_names[0]
        raw = pd.read_excel(path, sheet_name=sheet_name, header=None)
    except Exception:
        return pd.DataFrame()
    if raw is None or raw.empty:
        return pd.DataFrame()

    header_row = list(raw.iloc[0])
    col_idx = _build_col_idx(header_row)
    rows_iter = (tuple(r) for r in raw.iloc[1:].itertuples(index=False, name=None))

    def gi(*names):
        for n in names:
            if n in col_idx:
                return col_idx[n]
        return None

    idx_cabang = gi("CABANG")
    idx_tgl = gi("TGL FAKTUR", "TANGGAL")
    idx_kategori = gi("KATEGORI PENJUALAN")
    idx_total = gi("TOTAL HARGA")
    idx_qty = gi("QTY")
    idx_gp = gi("GROSS PROFIT")
    idx_pilar = _find_pilar_column_index(col_idx)
    idx_penjual = _find_penjual_column_index(col_idx)

    cabang_fallback = cabang_hint or branch_from_filename(os.path.basename(path)) or branch_from_sheetname(sheet_name)

    records = []
    for row in rows_iter:
        if row is None:
            continue
        total = _to_float_or_none(row[idx_total]) if idx_total is not None and idx_total < len(row) else None
        if total is None:
            continue
        cabang = _nan_to_none(row[idx_cabang]) if idx_cabang is not None and idx_cabang < len(row) else None
        cabang = str(cabang).strip().upper() if cabang else cabang_fallback
        tgl = to_date(row[idx_tgl]) if idx_tgl is not None and idx_tgl < len(row) else None
        kategori_raw = _nan_to_none(row[idx_kategori]) if idx_kategori is not None and idx_kategori < len(row) else None
        pilar_raw = _nan_to_none(row[idx_pilar]) if idx_pilar is not None and idx_pilar < len(row) else None
        penjual_raw = _nan_to_none(row[idx_penjual]) if idx_penjual is not None and idx_penjual < len(row) else None
        qty = _to_float_or_none(row[idx_qty]) if idx_qty is not None and idx_qty < len(row) else 0.0
        gp = _to_float_or_none(row[idx_gp]) if idx_gp is not None and idx_gp < len(row) else 0.0
        records.append({
            "Cabang": cabang,
            "Tanggal": tgl,
            "Kategori": classify_kategori(kategori_raw),
            "Omset": total,
            "Qty": qty or 0.0,
            "GrossProfit": gp or 0.0,
            "Pilar": classify_pilar(pilar_raw),
            "NamaPenjual": str(penjual_raw).strip() if penjual_raw else "",
            "PenjualKelompok": classify_penjual_kelompok(penjual_raw),
        })
    return pd.DataFrame(records)


def _load_master_sheet(path: str) -> pd.DataFrame:
    """Loader untuk format lama: master file dengan sheet 'Faktur Penjualan'.
    Pakai pandas.read_excel (bukan openpyxl read_only) supaya tahan terhadap
    file dengan metadata dimensi sheet yang tidak akurat."""
    try:
        xls = pd.ExcelFile(path)
    except Exception:
        return pd.DataFrame()
    all_sheets = xls.sheet_names
    sheet_candidates = [s for s in all_sheets if "FAKTUR" in s.upper() or s.upper().startswith("FP ")]
    if MAIN_SHEET_NAME in all_sheets:
        sheet_candidates = [MAIN_SHEET_NAME] + [s for s in sheet_candidates if s != MAIN_SHEET_NAME]
    if not sheet_candidates:
        sheet_candidates = [all_sheets[0]]

    frames = []
    for sheet_name in sheet_candidates:
        try:
            raw = pd.read_excel(path, sheet_name=sheet_name, header=None)
        except Exception:
            continue
        if raw is None or raw.empty:
            continue
        header_row = list(raw.iloc[0])
        col_idx = _build_col_idx(header_row)
        if "TOTAL HARGA" not in col_idx:
            continue
        rows_iter = (tuple(r) for r in raw.iloc[1:].itertuples(index=False, name=None))

        def gi(*names):
            for n in names:
                if n in col_idx:
                    return col_idx[n]
            return None

        idx_cabang = gi("CABANG")
        idx_tgl = gi("TGL FAKTUR", "TANGGAL")
        idx_kategori = gi("KATEGORI PENJUALAN")
        idx_total = gi("TOTAL HARGA")
        idx_qty = gi("QTY")
        idx_gp = gi("GROSS PROFIT")
        idx_pilar = _find_pilar_column_index(col_idx)
        idx_penjual = _find_penjual_column_index(col_idx)
        cabang_fallback = branch_from_sheetname(sheet_name)

        for row in rows_iter:
            if row is None:
                continue
            total = _to_float_or_none(row[idx_total]) if idx_total is not None and idx_total < len(row) else None
            if total is None:
                continue
            cabang = _nan_to_none(row[idx_cabang]) if idx_cabang is not None and idx_cabang < len(row) else None
            cabang = str(cabang).strip().upper() if cabang else cabang_fallback
            tgl = to_date(row[idx_tgl]) if idx_tgl is not None and idx_tgl < len(row) else None
            kategori_raw = _nan_to_none(row[idx_kategori]) if idx_kategori is not None and idx_kategori < len(row) else None
            pilar_raw = _nan_to_none(row[idx_pilar]) if idx_pilar is not None and idx_pilar < len(row) else None
            penjual_raw = _nan_to_none(row[idx_penjual]) if idx_penjual is not None and idx_penjual < len(row) else None
            qty = _to_float_or_none(row[idx_qty]) if idx_qty is not None and idx_qty < len(row) else 0.0
            gp = _to_float_or_none(row[idx_gp]) if idx_gp is not None and idx_gp < len(row) else 0.0
            frames.append({
                "Cabang": cabang,
                "Tanggal": tgl,
                "Kategori": classify_kategori(kategori_raw),
                "Omset": total,
                "Qty": qty or 0.0,
                "GrossProfit": gp or 0.0,
                "Pilar": classify_pilar(pilar_raw),
                "NamaPenjual": str(penjual_raw).strip() if penjual_raw else "",
                "PenjualKelompok": classify_penjual_kelompok(penjual_raw),
            })
    return pd.DataFrame(frames)


def _looks_like_ads_export(path: str) -> bool:
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        ws = wb[wb.sheetnames[0]]
        header = next(ws.iter_rows(values_only=True), None)
        wb.close()
        if not header:
            return False
        up = [str(h).upper() for h in header if h]
        return any("CAMPAIGN" in h or "AD SET" in h or "AMOUNT SPENT" in h for h in up)
    except Exception:
        return False


def _detect_main_file_kind(path: str):
    """Deteksi apakah file adalah format baru (per-cabang) atau lama (master)."""
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        names_up = [s.upper() for s in wb.sheetnames]
        wb.close()
        if any("RINCIAN FAKTUR" in n or "RINCIAN PENGIRIMAN" in n for n in names_up):
            if any("RINCIAN PENGIRIMAN" in n for n in names_up) and not any("RINCIAN FAKTUR" in n for n in names_up):
                return "walkin"
            return "faktur"
        if "FAKTUR PENJUALAN" in names_up or any(n.startswith("FP ") for n in names_up):
            return "master"
    except Exception:
        pass
    return None


def load_main_data(path: str) -> pd.DataFrame:
    kind = _detect_main_file_kind(path)
    if kind == "faktur":
        return _load_faktur_sheet(path)
    if kind == "master":
        return _load_master_sheet(path)
    try:
        return _load_faktur_sheet(path)
    except Exception:
        return pd.DataFrame()


def _dedupe_main_files():
    """Hapus file Omset duplikat/basi per cabang, sisakan yang timestamp-nya terbaru."""
    if not os.path.isdir(MAIN_DATA_DIR):
        return
    files = [f for f in os.listdir(MAIN_DATA_DIR) if f.lower().endswith((".xlsx", ".xls"))]
    groups = {}
    for f in files:
        branch = branch_from_filename(f) or "UNKNOWN"
        ts = _extract_filename_timestamp(f) or ""
        groups.setdefault(branch, []).append((ts, f))
    for branch, items in groups.items():
        if branch == "UNKNOWN" or len(items) <= 1:
            continue
        items.sort(key=lambda x: x[0])
        stale = items[:-1]
        for ts, fname in stale:
            fpath = os.path.join(MAIN_DATA_DIR, fname)
            try:
                os.remove(fpath)
                if _GH_ENABLED:
                    github_delete_file(f"data/main/{fname}", "auto-dedupe stale branch file")
            except Exception:
                pass


def load_all_main_data() -> pd.DataFrame:
    if not os.path.isdir(MAIN_DATA_DIR):
        return pd.DataFrame()
    frames = []
    for fname in sorted(os.listdir(MAIN_DATA_DIR)):
        if not fname.lower().endswith((".xlsx", ".xls")):
            continue
        fpath = os.path.join(MAIN_DATA_DIR, fname)
        try:
            df = load_main_data(fpath)
            if not df.empty:
                frames.append(df)
        except Exception:
            continue
    if not frames:
        return pd.DataFrame()
    combined = pd.concat(frames, ignore_index=True)
    combined["Tahun"] = combined["Tanggal"].apply(lambda d: d.year if d else None)
    combined["Bulan"] = combined["Tanggal"].apply(lambda d: d.month if d else None)
    combined = combined.dropna(subset=["Cabang", "Tanggal"])
    return combined


# ========================= Loader Data Iklan (Meta Ads) =========================

_ADS_REQUIRED_COLS = ["Campaign name", "Amount spent", "Reach", "Impressions"]


def _dedupe_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Buang kolom dengan nama duplikat (sisakan kemunculan pertama), supaya
    pd.concat tidak gagal dengan pandas.errors.InvalidIndexError."""
    if df is None or df.empty:
        return df
    return df.loc[:, ~df.columns.duplicated()]


def load_ads_data(path: str) -> pd.DataFrame:
    try:
        df = pd.read_excel(path)
    except Exception:
        return pd.DataFrame()
    df.columns = [str(c).strip() for c in df.columns]
    df = _dedupe_columns(df)
    if not any("amount spent" in c.lower() or c.lower() == "campaign name" for c in df.columns):
        return pd.DataFrame()
    branch = branch_from_filename(os.path.basename(path))
    df["Cabang"] = branch

    rename_map = {}
    used_targets = set()

    def _claim(col, target):
        if target in used_targets:
            return
        rename_map[col] = target
        used_targets.add(target)

    for c in df.columns:
        cl = c.lower()
        if "amount spent" in cl:
            _claim(c, "AmountSpent")
        elif cl == "reach":
            _claim(c, "Reach")
        elif cl == "impressions":
            _claim(c, "Impressions")
        elif "link click" in cl or cl == "clicks (all)":
            _claim(c, "Clicks")
        elif "cpm" in cl:
            _claim(c, "CPM")
        elif "cpc" in cl:
            _claim(c, "CPC")
        elif "ctr" in cl:
            _claim(c, "CTR")
        elif cl == "results":
            _claim(c, "Results")
        elif "messaging conversation" in cl and "cost per" not in cl:
            _claim(c, "Results")
        elif cl == "campaign name":
            _claim(c, "CampaignName")
    df = df.rename(columns=rename_map)
    df = _dedupe_columns(df)
    return df


def load_all_ads_data() -> pd.DataFrame:
    if not os.path.isdir(ADS_DATA_DIR):
        return pd.DataFrame()
    frames = []
    for fname in sorted(os.listdir(ADS_DATA_DIR)):
        if not fname.lower().endswith((".xlsx", ".xls", ".csv")):
            continue
        fpath = os.path.join(ADS_DATA_DIR, fname)
        try:
            df = load_ads_data(fpath)
            if df is not None and not df.empty:
                frames.append(_dedupe_columns(df))
        except Exception:
            continue
    if not frames:
        return pd.DataFrame()
    try:
        return pd.concat(frames, ignore_index=True, sort=False)
    except Exception:
        keep_cols = ["Cabang", "CampaignName", "AmountSpent", "Reach", "Impressions",
                     "Clicks", "CPM", "CPC", "CTR", "Results"]
        cleaned = []
        for f in frames:
            f = _dedupe_columns(f)
            cols_present = [c for c in keep_cols if c in f.columns]
            cleaned.append(f[cols_present])
        return pd.concat(cleaned, ignore_index=True, sort=False)


def aggregate_ads_by_branch(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty or "Cabang" not in df.columns:
        return pd.DataFrame()
    agg_map = {}
    for c in ["AmountSpent", "Reach", "Impressions", "Clicks", "Results"]:
        if c in df.columns:
            agg_map[c] = "sum"
    if not agg_map:
        return pd.DataFrame()
    g = df.groupby("Cabang").agg(agg_map).reset_index()
    return g


def _content_diagnosis(row):
    spend = row.get("AmountSpent", 0) or 0
    results = row.get("Results", 0) or 0
    if spend > 0 and results == 0:
        return "Spend ada tapi belum ada hasil — cek relevansi konten & targeting."
    if results > 0 and spend / max(results, 1) > 50000:
        return "Cost per result tinggi — evaluasi kreatif & audiens."
    return "Performa dalam batas wajar."


def generate_ads_insights(agg_df: pd.DataFrame) -> list:
    insights = []
    if agg_df.empty:
        return insights
    for _, r in agg_df.iterrows():
        diag = _content_diagnosis(r)
        if "cek" in diag.lower() or "evaluasi" in diag.lower():
            insights.append({
                "level": "warn",
                "category": r["Cabang"],
                "title": f"{r['Cabang']}: {diag}",
                "problem": f"Amount Spent {format_rupiah(r.get('AmountSpent', 0))}, Results {format_number(r.get('Results', 0))}.",
                "online": ["Uji ulang kreatif iklan (gambar/video/copy)", "Perbaiki targeting audiens berdasar data insight"],
                "offline": ["Selaraskan promo iklan dengan promo di toko"],
            })
    return insights


def render_insight_card(ins: dict):
    level_color = {"bad": "#dc2626", "warn": "#d97706", "good": "#16a34a"}.get(ins.get("level"), "#6b7280")
    online_html = "".join(f"<li>{x}</li>" for x in ins.get("online", []))
    offline_html = "".join(f"<li>{x}</li>" for x in ins.get("offline", []))
    st.markdown(
        f"""
        <div style="border-left:4px solid {level_color};background:#f9fafb;padding:10px 14px;border-radius:6px;margin-bottom:10px;">
            <div style="font-weight:700;color:{level_color};">{ins.get('title','')}</div>
            <div style="font-size:0.9em;color:#374151;margin:4px 0;">{ins.get('problem','')}</div>
            <div style="display:flex;gap:20px;font-size:0.85em;">
                <div><b>Online:</b><ul style="margin:2px 0;">{online_html}</ul></div>
                <div><b>Offline:</b><ul style="margin:2px 0;">{offline_html}</ul></div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


# ========================= Sales Insight Engine =========================

_SALES_TOTAL_LABELS = {"Service": "Omset Service", "Gadget & Aksesoris": "Omset Gadget & Aksesoris"}
_CATEGORY_ACTION_PLANS = {
    "Service": {
        "online": ["Promosikan layanan service via media sosial & Google Business", "Tawarkan booking service online"],
        "offline": ["Tingkatkan kualitas layanan & waktu pengerjaan", "Promo diskon service di jam sepi"],
    },
    "Gadget & Aksesoris": {
        "online": ["Optimalkan katalog produk di marketplace & media sosial", "Jalankan campaign promo gadget/aksesoris"],
        "offline": ["Perbarui display produk di toko", "Bundling promo gadget + aksesoris"],
    },
    "Marketing Corporate": {
        "online": ["Follow-up leads corporate via WhatsApp Business/LinkedIn"],
        "offline": ["Kunjungan langsung ke calon klien corporate"],
    },
}


def generate_sales_insights(df_branch_month: pd.DataFrame, kategori: str) -> list:
    insights = []
    if df_branch_month.empty or len(df_branch_month) < 2:
        return insights
    df_sorted = df_branch_month.sort_values(["Tahun", "Bulan"])
    if len(df_sorted) < 2:
        return insights
    last, prev = df_sorted.iloc[-1], df_sorted.iloc[-2]
    delta = last["Total"] - prev["Total"]
    if prev["Total"] and delta / prev["Total"] < -0.1:
        plan = _CATEGORY_ACTION_PLANS.get(kategori, {"online": [], "offline": []})
        insights.append({
            "level": "bad",
            "category": kategori,
            "title": f"{_SALES_TOTAL_LABELS.get(kategori, kategori)} turun {format_percent(abs(delta / prev['Total']))} dibanding bulan lalu",
            "problem": f"Dari {format_rupiah(prev['Total'])} menjadi {format_rupiah(last['Total'])}.",
            "online": plan["online"], "offline": plan["offline"],
        })
    return insights


def generate_all_sales_insights(df: pd.DataFrame) -> list:
    insights = []
    if df.empty:
        return insights
    for kategori in ["Service", "Gadget & Aksesoris"]:
        sub = df[df["Kategori"] == kategori]
        if sub.empty:
            continue
        g = sub.groupby(["Tahun", "Bulan"])["Omset"].sum().reset_index().rename(columns={"Omset": "Total"})
        insights.extend(generate_sales_insights(g, kategori))
    return insights


def _render_online_offline_html(online: list, offline: list) -> str:
    online_html = "".join(f"<li>{x}</li>" for x in online)
    offline_html = "".join(f"<li>{x}</li>" for x in offline)
    return f"""
    <div style="display:flex;gap:20px;font-size:0.85em;">
        <div><b>Online:</b><ul style="margin:2px 0;">{online_html}</ul></div>
        <div><b>Offline:</b><ul style="margin:2px 0;">{offline_html}</ul></div>
    </div>
    """


def render_structured_insight_card(ins: dict):
    level_color = {"bad": "#dc2626", "warn": "#d97706", "good": "#16a34a"}.get(ins.get("level"), "#6b7280")
    st.markdown(
        f"""
        <div style="border-left:4px solid {level_color};background:#f9fafb;padding:10px 14px;border-radius:6px;margin-bottom:10px;">
            <div style="font-weight:700;color:{level_color};">{ins.get('title','')}</div>
            <div style="font-size:0.9em;color:#374151;margin:4px 0;">{ins.get('problem','')}</div>
            {_render_online_offline_html(ins.get('online', []), ins.get('offline', []))}
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_kpi_card_text(label: str, value: str, color: str = "#1d4ed8", icon: str = "", sub: str = ""):
    sub_html = f'<div style="font-size:0.75em;color:#6b7280;">{sub}</div>' if sub else ""
    return f"""
    <div style="border:2px solid {color};border-radius:10px;padding:12px 16px;text-align:center;background:white;">
        <div style="font-size:1.6em;">{icon}</div>
        <div style="font-size:0.85em;font-weight:700;color:{color};margin:4px 0;">{label}</div>
        <div style="font-size:1.15em;font-weight:800;color:#111827;">{value}</div>
        {sub_html}
    </div>
    """


def render_kpi_card(label: str, value: str, color: str = "#1d4ed8", icon: str = ""):
    return render_kpi_card_text(label, value, color, icon)


# ========================= Loader Data Walk-in =========================

def load_walkin_data(path: str) -> pd.DataFrame:
    # Pakai pandas.read_excel (bukan openpyxl read_only) karena sejumlah file
    # export "Rincian Pengiriman Pesanan" MFlash punya metadata dimensi sheet yang
    # tidak akurat sehingga openpyxl read_only gagal mendeteksi baris data.
    try:
        sheet_name = None
        try:
            xls = pd.ExcelFile(path)
            for s in xls.sheet_names:
                if "PENGIRIMAN" in s.upper():
                    sheet_name = s
                    break
            if sheet_name is None:
                sheet_name = xls.sheet_names[0]
        except Exception:
            sheet_name = 0
        raw = pd.read_excel(path, sheet_name=sheet_name)
    except Exception:
        return pd.DataFrame()
    if raw is None or raw.empty:
        return pd.DataFrame()
    raw.columns = [str(c).strip().upper() for c in raw.columns]
    raw = _dedupe_columns(raw)

    def gi(*names):
        for n in names:
            if n in raw.columns:
                return n
        for n in names:
            for c in raw.columns:
                if n in c:
                    return c
        return None

    col_tgl = gi("TGL PENGIRIMAN", "TANGGAL PENGIRIMAN", "TANGGAL")
    col_nomor = gi("NOMOR PENGIRIMAN PESANAN", "NOMOR PENGIRIMAN", "NO PENGIRIMAN", "NO. PENGIRIMAN")
    col_cabang = gi("CABANG")
    if col_tgl is None or col_nomor is None:
        return pd.DataFrame()

    cabang_fallback = branch_from_filename(os.path.basename(path)) or (
        branch_from_sheetname(sheet_name) if isinstance(sheet_name, str) else None
    )
    records = []
    for _, row in raw.iterrows():
        tgl = to_date(row[col_tgl])
        nomor = row[col_nomor]
        if tgl is None or nomor is None or (isinstance(nomor, float) and pd.isna(nomor)) or str(nomor).strip() == "":
            continue
        cabang = row[col_cabang] if col_cabang is not None else None
        cabang = str(cabang).strip().upper() if cabang and str(cabang).strip() else cabang_fallback
        records.append({"Cabang": cabang, "Tanggal": tgl, "NomorPengiriman": str(nomor).strip()})
    df = pd.DataFrame(records)
    if df.empty:
        return df
    df["Tahun"] = df["Tanggal"].apply(lambda d: d.year)
    df["Bulan"] = df["Tanggal"].apply(lambda d: d.month)
    return df


def load_all_walkin_data() -> pd.DataFrame:
    if not os.path.isdir(WALKIN_DATA_DIR):
        return pd.DataFrame()
    frames = []
    for fname in sorted(os.listdir(WALKIN_DATA_DIR)):
        if not fname.lower().endswith((".xlsx", ".xls")):
            continue
        fpath = os.path.join(WALKIN_DATA_DIR, fname)
        try:
            df = load_walkin_data(fpath)
            if not df.empty:
                frames.append(df)
        except Exception:
            continue
    if not frames:
        return pd.DataFrame()
    combined = pd.concat(frames, ignore_index=True)
    combined = combined.dropna(subset=["Cabang", "Tanggal"])
    combined = combined.drop_duplicates(subset=["Cabang", "NomorPengiriman"])
    return combined


def aggregate_walkin_monthly(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return pd.DataFrame(columns=["Cabang", "Tahun", "Bulan", "TotalWalkin", "RataRataPerHari"])
    rows = []
    for (cabang, tahun, bulan), g in df.groupby(["Cabang", "Tahun", "Bulan"]):
        total = g["NomorPengiriman"].nunique()
        hari_dalam_bulan = calendar.monthrange(int(tahun), int(bulan))[1]
        rata2 = total / hari_dalam_bulan if hari_dalam_bulan else 0.0
        rows.append({"Cabang": cabang, "Tahun": tahun, "Bulan": bulan, "TotalWalkin": int(total), "RataRataPerHari": rata2})
    return pd.DataFrame(rows)


def aggregate_walkin_current_period(df: pd.DataFrame, tanggal_acuan: date, selected_branches=None) -> pd.DataFrame:
    """Total & rata-rata Walk-in per cabang KUMULATIF dari awal kuartal berjalan
    (1 Jan/Apr/Jul/Okt) sampai Tanggal Acuan — SAMA seperti kolom 'S/D HARI INI' di
    Scoreboard, supaya konsisten dengan seluruh dashboard dan sesuai kebutuhan bisnis
    (total walk-in sejak awal periode s/d tanggal terakhir data)."""
    cols = ["Cabang", "TotalWalkin", "HariEfektif", "RataRataPerHari"]
    if df.empty or tanggal_acuan is None:
        return pd.DataFrame(columns=cols)
    start, end, total_hari, hari_berjalan, sisa_hari = _quarter_bounds(tanggal_acuan)
    mask = (df["Tanggal"] >= start) & (df["Tanggal"] <= tanggal_acuan)
    d = df[mask]
    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]
    if d.empty:
        return pd.DataFrame(columns=cols)
    hari_efektif = hari_berjalan
    rows = []
    for cabang, g in d.groupby("Cabang"):
        total = g["NomorPengiriman"].nunique()
        rata2 = (total / hari_efektif) if hari_efektif else 0.0
        rows.append({
            "Cabang": cabang, "TotalWalkin": int(total), "HariEfektif": int(hari_efektif), "RataRataPerHari": rata2,
        })
    return pd.DataFrame(rows).reset_index(drop=True)


def _walkin_ordered(d: pd.DataFrame) -> pd.DataFrame:
    if d.empty:
        return d
    d = d.copy()
    d["_order"] = d["Cabang"].apply(lambda b: _BRANCH_RANK.get(str(b).upper(), 999))
    return d.sort_values("_order").drop(columns="_order").reset_index(drop=True)


def _walkin_overall_avg(d: pd.DataFrame) -> float:
    if d.empty:
        return 0.0
    return float(d["RataRataPerHari"].mean())


def render_walkin_table_html(d: pd.DataFrame, periode_label: str = "") -> str:
    if d.empty:
        return "<i>Tidak ada data walk-in untuk periode ini.</i>"
    overall_avg = _walkin_overall_avg(d)
    header_label = f" — {periode_label}" if periode_label else ""
    rows_html = ""
    for _, r in d.iterrows():
        above = r["RataRataPerHari"] >= overall_avg
        bg = "#dcfce7" if above else "#fee2e2"
        badge = "🟢" if above else "🔴"
        rows_html += f"""<tr style="background:{bg};">
            <td style="padding:6px 10px;border:1px solid #d1d5db;font-weight:600;">{r['Cabang']}</td>
            <td style="padding:6px 10px;border:1px solid #d1d5db;text-align:right;">{format_number(r['TotalWalkin'])}</td>
            <td style="padding:6px 10px;border:1px solid #d1d5db;text-align:right;">{format_decimal(r['RataRataPerHari'])} {badge}</td>
        </tr>"""
    return f"""
    <div style="border:2px solid #0f766e;border-radius:10px;overflow:hidden;">
    <div style="background:#0f766e;color:white;padding:8px 12px;font-weight:700;">
        📋 Tabel Jumlah &amp; Rata-rata Walk-in per Cabang{header_label}
    </div>
    <table style="border-collapse:collapse;width:100%;font-size:0.92em;">
        <thead><tr style="background:#f0fdfa;">
            <th style="padding:6px 10px;border:1px solid #d1d5db;text-align:left;">Cabang</th>
            <th style="padding:6px 10px;border:1px solid #d1d5db;">Total Walk-in</th>
            <th style="padding:6px 10px;border:1px solid #d1d5db;">Rata-rata/Hari</th>
        </tr></thead>
        <tbody>{rows_html}</tbody>
    </table>
    </div>
    """


def generate_walkin_table_image(d: pd.DataFrame, periode_label: str = "") -> bytes:
    title0 = "Tabel Jumlah & Rata-rata Walk-in per Cabang"
    if periode_label:
        title0 += f" — {periode_label}"
    if d.empty:
        fig, ax = plt.subplots(figsize=(7, 2.2))
        ax.axis("off")
        ax.set_title(title0, fontsize=12, fontweight="bold", color="#0f766e", pad=14)
        ax.text(0.5, 0.5, "Tidak ada data walk-in untuk periode ini.", ha="center", va="center", fontsize=10, color="#6b7280")
        buf = io.BytesIO()
        fig.savefig(buf, format="jpg", dpi=200, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    overall_avg = _walkin_overall_avg(d)
    n = len(d)
    fig_h = max(1.5, 0.5 * n + 1.2)
    fig, ax = plt.subplots(figsize=(7, fig_h))
    ax.axis("off")
    ax.set_title(title0, fontsize=12, fontweight="bold", color="#0f766e", pad=14)

    col_labels = ["Cabang", "Total Walk-in", "Rata-rata/Hari"]
    cell_text = [[r["Cabang"], format_number(r["TotalWalkin"]), format_decimal(r["RataRataPerHari"])] for _, r in d.iterrows()]
    tbl = ax.table(cellText=cell_text, colLabels=col_labels, loc="center", cellLoc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9)
    tbl.scale(1, 1.6)
    for (row, col), cell in tbl.get_celld().items():
        cell.set_edgecolor("#d1d5db")
        if row == 0:
            cell.set_facecolor("#0f766e")
            cell.set_text_props(color="white", fontweight="bold")
        else:
            above = d.iloc[row - 1]["RataRataPerHari"] >= overall_avg
            cell.set_facecolor("#dcfce7" if above else "#fee2e2")
    buf = io.BytesIO()
    fig.tight_layout()
    fig.savefig(buf, format="jpg", dpi=200, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def generate_walkin_table_pdf(d: pd.DataFrame, periode_label: str = "") -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=1.5 * cm, bottomMargin=1.5 * cm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("title", parent=styles["Title"], fontSize=14, textColor=rl_colors.HexColor("#0f766e"))
    elements = [Paragraph(f"Tabel Jumlah & Rata-rata Walk-in per Cabang{' — ' + periode_label if periode_label else ''}", title_style), Spacer(1, 12)]

    if d.empty:
        elements.append(Paragraph("Tidak ada data walk-in untuk periode ini.", styles["Normal"]))
    else:
        overall_avg = _walkin_overall_avg(d)
        data = [["Cabang", "Total Walk-in", "Rata-rata/Hari"]]
        row_colors = []
        for _, r in d.iterrows():
            data.append([r["Cabang"], format_number(r["TotalWalkin"]), format_decimal(r["RataRataPerHari"])])
            row_colors.append(r["RataRataPerHari"] >= overall_avg)
        table = Table(data, colWidths=[6 * cm, 5 * cm, 5 * cm])
        style_cmds = [
            ("BACKGROUND", (0, 0), (-1, 0), rl_colors.HexColor("#0f766e")),
            ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.HexColor("#d1d5db")),
            ("ALIGN", (1, 0), (-1, -1), "CENTER"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]
        for i, above in enumerate(row_colors, start=1):
            bg = rl_colors.HexColor("#dcfce7") if above else rl_colors.HexColor("#fee2e2")
            style_cmds.append(("BACKGROUND", (0, i), (-1, i), bg))
        table.setStyle(TableStyle(style_cmds))
        elements.append(table)

    doc.build(elements)
    buf.seek(0)
    return buf.read()


_WALKIN_ACTION_PLAN = {
    "online": [
        "Promosi lokasi via Google Maps/Instagram Ads radius sekitar cabang",
        "Kampanye promo walk-in (diskon cek gratis, hari tertentu)",
    ],
    "offline": [
        "Spanduk/banner promo di depan toko & area sekitar",
        "Kerja sama dengan warga/komunitas sekitar (RT/RW, kampus, kantor)",
    ],
}


def generate_walkin_marketing_insights(d: pd.DataFrame) -> list:
    insights = []
    if d.empty:
        return insights
    overall_avg = _walkin_overall_avg(d)
    if overall_avg <= 0:
        return insights
    for _, r in d.iterrows():
        if r["RataRataPerHari"] < 0.85 * overall_avg:
            insights.append({
                "level": "bad",
                "category": r["Cabang"],
                "title": f"{r['Cabang']}: Walk-in di bawah rata-rata cabang lain",
                "problem": f"Rata-rata {format_decimal(r['RataRataPerHari'])}/hari, sementara rata-rata seluruh cabang {format_decimal(overall_avg)}/hari.",
                "online": _WALKIN_ACTION_PLAN["online"], "offline": _WALKIN_ACTION_PLAN["offline"],
            })
    return insights


def generate_walkin_insights(walkin_agg: pd.DataFrame) -> list:
    """Insight tren bulan-ke-bulan (dibanding bulan lalu), pakai riwayat multi-bulan."""
    insights = []
    if walkin_agg.empty:
        return insights
    for cabang, g in walkin_agg.groupby("Cabang"):
        g = g.sort_values(["Tahun", "Bulan"])
        if len(g) < 2:
            continue
        last, prev = g.iloc[-1], g.iloc[-2]
        if prev["RataRataPerHari"] and (last["RataRataPerHari"] - prev["RataRataPerHari"]) / prev["RataRataPerHari"] < -0.15:
            insights.append({
                "level": "warn",
                "category": cabang,
                "title": f"{cabang}: Walk-in turun dibanding bulan lalu",
                "problem": f"Rata-rata/hari dari {format_decimal(prev['RataRataPerHari'])} menjadi {format_decimal(last['RataRataPerHari'])}.",
                "online": _WALKIN_ACTION_PLAN["online"], "offline": _WALKIN_ACTION_PLAN["offline"],
            })
    return insights


# ========================= Corporate & Target Loaders ==========================

def load_corporate_data(path: str) -> pd.DataFrame:
    try:
        df = pd.read_excel(path)
    except Exception:
        return pd.DataFrame()
    df.columns = [str(c).strip() for c in df.columns]
    df = _dedupe_columns(df)
    return df


def make_corporate_template() -> bytes:
    df = pd.DataFrame({"Cabang": BRANCH_ORDER, "Nama Marketing": ["" for _ in BRANCH_ORDER], "Omset Corporate": [0 for _ in BRANCH_ORDER]})
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Template Corporate")
    buf.seek(0)
    return buf.read()


def load_target_data(path: str) -> pd.DataFrame:
    try:
        df = pd.read_excel(path)
    except Exception:
        return pd.DataFrame()
    df.columns = [str(c).strip() for c in df.columns]
    df = _dedupe_columns(df)
    return df


def make_target_template() -> bytes:
    rows = []
    for kat in ["Omset All", "Service", "Gadget & Aksesoris"]:
        for b in BRANCH_ORDER:
            rows.append({"Kategori": kat, "Cabang": b, "Target Omset": 0})
    df = pd.DataFrame(rows)
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Template Target")
    buf.seek(0)
    return buf.read()


# ========================= Scoreboard Core ==========================

SCOREBOARD_KATEGORI = ["Omset All", "Service", "Gadget & Aksesoris"]
_SCOREBOARD_KATEGORI_LABEL = {
    "Omset All": "SCOREBOARD OMSET ALL",
    "Service": "SCOREBOARD OMSET SERVICE",
    "Gadget & Aksesoris": "SCOREBOARD OMSET GADGET & AKSESORIS",
}
MONEY_COLS = ["OmsetSamurai", "OmsetHarian", "ExpectedValue", "HariIni", "SdHariIni",
              "GapHariIni", "TotalGap", "KejarPerhari", "PeriodeBulanLalu", "PeriodeBulanIni"]


def _quarter_bounds(d: date):
    """(start, end, total_hari, hari_berjalan, sisa_hari) untuk kuartal kalender yang berisi d."""
    if d is None:
        d = date.today()
    q = (d.month - 1) // 3
    start_month = q * 3 + 1
    start = date(d.year, start_month, 1)
    end_month = start_month + 2
    end_year = d.year
    if end_month > 12:
        end_month -= 12
        end_year += 1
    last_day = calendar.monthrange(end_year, end_month)[1]
    end = date(end_year, end_month, last_day)
    total_hari = (end - start).days + 1
    hari_berjalan = (min(d, end) - start).days + 1
    hari_berjalan = max(1, hari_berjalan)
    sisa_hari = max(0, total_hari - hari_berjalan)
    return start, end, total_hari, hari_berjalan, sisa_hari


def pencapaian_color(pct: float) -> str:
    if pct is None:
        return "#9ca3af"
    if pct >= 1.0:
        return "#16a34a"
    if pct >= 0.85:
        return "#f59e0b"
    return "#dc2626"


def build_scoreboard(df_main: pd.DataFrame, tanggal_acuan: date, target_map: dict, kategori: str, selected_branches=None) -> pd.DataFrame:
    """Bangun tabel scoreboard untuk satu kategori, formula sama persis dengan sheet 'Data Periode' Excel."""
    cols = ["Cabang", "OmsetSamurai", "OmsetHarian", "TotalHari", "HariIni", "HariBerjalan", "SisaHari",
            "ExpectedValue", "SdHariIni", "PctPencapaian", "GapHariIni", "TotalGap", "KejarPerhari",
            "PeriodeBulanLalu", "PeriodeBulanIni"]
    if df_main is None or df_main.empty or tanggal_acuan is None:
        return pd.DataFrame(columns=cols)

    start, end, total_hari, hari_berjalan, sisa_hari = _quarter_bounds(tanggal_acuan)

    d = df_main.copy()
    if "Tanggal" not in d.columns:
        return pd.DataFrame(columns=cols)
    d = d[(d["Tanggal"] >= start) & (d["Tanggal"] <= tanggal_acuan)]

    if kategori == "Service":
        d = d[d["Kategori"] == "Service"]
    elif kategori == "Gadget & Aksesoris":
        d = d[d["Kategori"] == "Gadget & Aksesoris"]

    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]

    branches = selected_branches if selected_branches else BRANCH_ORDER
    branch_target = target_map.get(kategori, {}) if target_map else {}

    prev_month_start, prev_month_end = _prev_month_bounds(tanggal_acuan)
    cur_month_start = date(tanggal_acuan.year, tanggal_acuan.month, 1)

    rows = []
    for cabang in order_branches(branches):
        omset_samurai = float(branch_target.get(cabang, 0) or 0)
        g = d[d["Cabang"] == cabang]
        sd_hari_ini = float(g["Omset"].sum()) if not g.empty else 0.0

        omset_harian = omset_samurai / total_hari if total_hari else 0.0
        expected_value = omset_harian * hari_berjalan
        pct = (sd_hari_ini / expected_value) if expected_value else None
        gap_hari_ini = expected_value - sd_hari_ini
        total_gap = omset_samurai - sd_hari_ini
        kejar_perhari = (total_gap / sisa_hari) if sisa_hari else 0.0

        g_all = df_main[df_main["Cabang"] == cabang]
        if kategori == "Service":
            g_all = g_all[g_all["Kategori"] == "Service"]
        elif kategori == "Gadget & Aksesoris":
            g_all = g_all[g_all["Kategori"] == "Gadget & Aksesoris"]

        prev_g = g_all[(g_all["Tanggal"] >= prev_month_start) & (g_all["Tanggal"] <= prev_month_end)]
        prev_days = (prev_month_end - prev_month_start).days + 1
        periode_bulan_lalu = (float(prev_g["Omset"].sum()) / prev_days) if prev_days else 0.0

        cur_g = g_all[(g_all["Tanggal"] >= cur_month_start) & (g_all["Tanggal"] <= tanggal_acuan)]
        cur_days = (tanggal_acuan - cur_month_start).days + 1
        periode_bulan_ini = (float(cur_g["Omset"].sum()) / cur_days) if cur_days else 0.0

        rows.append({
            "Cabang": cabang, "OmsetSamurai": omset_samurai, "OmsetHarian": omset_harian,
            "TotalHari": total_hari, "HariIni": hari_berjalan, "HariBerjalan": hari_berjalan, "SisaHari": sisa_hari,
            "ExpectedValue": expected_value, "SdHariIni": sd_hari_ini, "PctPencapaian": pct,
            "GapHariIni": gap_hari_ini, "TotalGap": total_gap, "KejarPerhari": kejar_perhari,
            "PeriodeBulanLalu": periode_bulan_lalu, "PeriodeBulanIni": periode_bulan_ini,
        })
    return pd.DataFrame(rows)


def _prev_month_bounds(d: date):
    if d.month == 1:
        y, m = d.year - 1, 12
    else:
        y, m = d.year, d.month - 1
    start = date(y, m, 1)
    end = date(y, m, calendar.monthrange(y, m)[1])
    return start, end


def _finalize_scoreboard(df_sb: pd.DataFrame) -> pd.DataFrame:
    if df_sb.empty:
        return df_sb
    total_row = {"Cabang": "TOTAL"}
    for c in ["OmsetSamurai", "SdHariIni", "ExpectedValue", "GapHariIni", "TotalGap"]:
        total_row[c] = df_sb[c].sum()
    total_row["OmsetHarian"] = df_sb["OmsetHarian"].sum()
    total_row["TotalHari"] = df_sb["TotalHari"].iloc[0] if len(df_sb) else 0
    total_row["HariIni"] = df_sb["HariIni"].iloc[0] if len(df_sb) else 0
    total_row["SisaHari"] = df_sb["SisaHari"].iloc[0] if len(df_sb) else 0
    total_row["PctPencapaian"] = (total_row["SdHariIni"] / total_row["ExpectedValue"]) if total_row["ExpectedValue"] else None
    total_row["KejarPerhari"] = (total_row["TotalGap"] / total_row["SisaHari"]) if total_row["SisaHari"] else 0.0
    total_row["PeriodeBulanLalu"] = df_sb["PeriodeBulanLalu"].sum()
    total_row["PeriodeBulanIni"] = df_sb["PeriodeBulanIni"].sum()
    return pd.concat([df_sb, pd.DataFrame([total_row])], ignore_index=True)


# ========================= Target Auto-Extraction dari Sheet Scoreboard ==========================

_SECTION_MARKERS = {
    "Omset All": "SCOREBOARD OMSET ALL",
    "Service": "SCOREBOARD OMSET SERVICE",
    "Gadget & Aksesoris": "SCOREBOARD OMSET GADGET",
}


def _read_scoreboard_sections(path: str) -> dict:
    """Baca sheet 'Scoreboard' pada file Excel utama, ekstrak target Omset Samurai per kategori & cabang."""
    result = {k: {} for k in SCOREBOARD_KATEGORI}
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception:
        return result
    sheet_name = None
    for sn in wb.sheetnames:
        if "scoreboard" in sn.lower():
            sheet_name = sn
            break
    if not sheet_name:
        return result
    ws = wb[sheet_name]

    rows = list(ws.iter_rows(values_only=True))
    current_kat = None
    header_row_idx = None
    cabang_col = None
    samurai_col = None
    for i, row in enumerate(rows):
        row_text = " ".join([str(c) for c in row if c is not None]).upper()
        matched_kat = None
        for kat, marker in _SECTION_MARKERS.items():
            if marker in row_text:
                matched_kat = kat
                break
        if matched_kat:
            current_kat = matched_kat
            header_row_idx = None
            cabang_col = None
            samurai_col = None
            continue
        if current_kat is None:
            continue
        if header_row_idx is None:
            for j, cell in enumerate(row):
                if cell is None:
                    continue
                cl = str(cell).strip().upper()
                if cl == "CABANG":
                    cabang_col = j
                elif "SAMURAI" in cl or ("TARGET" in cl and "OMSET" in cl):
                    samurai_col = j
            if cabang_col is not None and samurai_col is not None:
                header_row_idx = i
            continue
        cabang_val = row[cabang_col] if cabang_col < len(row) else None
        samurai_val = row[samurai_col] if samurai_col < len(row) else None
        if cabang_val is None:
            continue
        cabang_str = str(cabang_val).strip().upper()
        if cabang_str in ("", "TOTAL", "GRAND TOTAL"):
            if cabang_str == "TOTAL":
                current_kat = None
            continue
        matched_branch = None
        for b in BRANCH_ORDER:
            if b == cabang_str or b in cabang_str or cabang_str in b:
                matched_branch = b
                break
        if matched_branch and samurai_val is not None:
            try:
                result[current_kat][matched_branch] = float(samurai_val)
            except (ValueError, TypeError):
                pass
    return result


def extract_scoreboard_target(path: str) -> dict:
    return _read_scoreboard_sections(path)


def extract_scoreboard_snapshot_date(path: str):
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception:
        return None
    sheet_name = None
    for sn in wb.sheetnames:
        if "periode" in sn.lower():
            sheet_name = sn
            break
    if not sheet_name:
        return None
    ws = wb[sheet_name]
    for row in ws.iter_rows(values_only=True):
        for cell in row:
            if isinstance(cell, datetime):
                return cell.date()
            if isinstance(cell, date):
                return cell
    return None


def extract_scoreboard_corporate(path: str) -> pd.DataFrame:
    """Stub: ekstraksi kontribusi corporate dari sheet Scoreboard (belum ada format standar)."""
    return pd.DataFrame(columns=["Cabang", "Nama Marketing", "Omset Corporate"])


# ========================= Scoreboard Render ==========================

def _fmt_scoreboard_cell(col: str, val) -> str:
    if val is None or (isinstance(val, float) and pd.isna(val)):
        return "-"
    if col == "PctPencapaian":
        return format_percent(val)
    if col in MONEY_COLS:
        return format_rupiah(val)
    if col in ("TotalHari", "HariIni", "HariBerjalan", "SisaHari"):
        return format_number(val)
    return str(val)


_SCOREBOARD_COL_ORDER = ["Cabang", "OmsetSamurai", "OmsetHarian", "TotalHari", "HariIni", "SisaHari",
                         "ExpectedValue", "SdHariIni", "PctPencapaian", "GapHariIni", "TotalGap", "KejarPerhari",
                         "PeriodeBulanLalu", "PeriodeBulanIni"]
_SCOREBOARD_GROUPS = {
    "Cabang": "Cabang", "OmsetSamurai": "Target (Omset Samurai)", "OmsetHarian": "Omset Harian",
    "TotalHari": "Total Hari", "HariIni": "Hari Berjalan", "SisaHari": "Sisa Hari",
    "ExpectedValue": "Expected Value", "SdHariIni": "S/D Hari Ini", "PctPencapaian": "% Pencapaian",
    "GapHariIni": "Gap Hari Ini", "TotalGap": "Total Gap", "KejarPerhari": "Kejar/Hari",
    "PeriodeBulanLalu": "Rata2 Bulan Lalu", "PeriodeBulanIni": "Rata2 Bulan Ini",
}


def render_scoreboard_html(df_sb: pd.DataFrame, kategori: str) -> str:
    if df_sb.empty:
        return "<i>Tidak ada data scoreboard untuk kategori ini.</i>"
    header_cells = "".join(f'<th style="padding:5px 8px;border:1px solid #d1d5db;white-space:nowrap;">{_SCOREBOARD_GROUPS.get(c, c)}</th>' for c in _SCOREBOARD_COL_ORDER)
    body_rows = ""
    for _, r in df_sb.iterrows():
        is_total = str(r["Cabang"]).upper() == "TOTAL"
        row_bg = "#f0fdfa" if is_total else "white"
        fw = "700" if is_total else "400"
        cells = ""
        for c in _SCOREBOARD_COL_ORDER:
            val = r.get(c)
            txt = _fmt_scoreboard_cell(c, val)
            align = "left" if c == "Cabang" else "right"
            extra_style = ""
            if c == "PctPencapaian" and val is not None and not (isinstance(val, float) and pd.isna(val)):
                extra_style = f"color:{pencapaian_color(val)};font-weight:700;"
            cells += f'<td style="padding:5px 8px;border:1px solid #d1d5db;text-align:{align};font-weight:{fw};{extra_style}">{txt}</td>'
        body_rows += f'<tr style="background:{row_bg};">{cells}</tr>'
    title = _SCOREBOARD_KATEGORI_LABEL.get(kategori, kategori)
    return f"""
    <div style="border:2px solid #0f766e;border-radius:10px;overflow-x:auto;margin-bottom:14px;">
    <div style="background:#0f766e;color:white;padding:8px 12px;font-weight:700;">🏆 {title}</div>
    <table style="border-collapse:collapse;width:100%;font-size:0.85em;">
        <thead><tr style="background:#ccfbf1;">{header_cells}</tr></thead>
        <tbody>{body_rows}</tbody>
    </table>
    </div>
    """


# ========================= Scoreboard Export (JPG/PDF) ==========================

def generate_scoreboard_table_image(df_sb: pd.DataFrame, kategori: str) -> bytes:
    title0 = _SCOREBOARD_KATEGORI_LABEL.get(kategori, kategori)
    if df_sb.empty:
        fig, ax = plt.subplots(figsize=(10, 2.2))
        ax.axis("off")
        ax.set_title(title0, fontsize=12, fontweight="bold", color="#0f766e", pad=14)
        ax.text(0.5, 0.5, "Tidak ada data untuk kategori ini.", ha="center", va="center", fontsize=10, color="#6b7280")
        buf = io.BytesIO()
        fig.savefig(buf, format="jpg", dpi=200, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    col_labels = [_SCOREBOARD_GROUPS.get(c, c) for c in _SCOREBOARD_COL_ORDER]
    cell_text = []
    for _, r in df_sb.iterrows():
        cell_text.append([_fmt_scoreboard_cell(c, r.get(c)) for c in _SCOREBOARD_COL_ORDER])

    n = len(df_sb)
    fig_w = max(12, 1.1 * len(_SCOREBOARD_COL_ORDER))
    fig_h = max(2.0, 0.45 * n + 1.4)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.axis("off")
    ax.set_title(title0, fontsize=13, fontweight="bold", color="#0f766e", pad=16)

    tbl = ax.table(cellText=cell_text, colLabels=col_labels, loc="center", cellLoc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(8)
    tbl.scale(1, 1.5)
    for (row, col), cell in tbl.get_celld().items():
        cell.set_edgecolor("#d1d5db")
        if row == 0:
            cell.set_facecolor("#0f766e")
            cell.set_text_props(color="white", fontweight="bold")
        else:
            is_total = str(df_sb.iloc[row - 1]["Cabang"]).upper() == "TOTAL"
            cell.set_facecolor("#f0fdfa" if is_total else "white")
            col_name = _SCOREBOARD_COL_ORDER[col]
            if col_name == "PctPencapaian":
                pct_val = df_sb.iloc[row - 1].get("PctPencapaian")
                if pct_val is not None and not (isinstance(pct_val, float) and pd.isna(pct_val)):
                    cell.set_text_props(color=pencapaian_color(pct_val), fontweight="bold")
    buf = io.BytesIO()
    fig.tight_layout()
    fig.savefig(buf, format="jpg", dpi=200, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def generate_scoreboard_pdf(sb_dict: dict) -> bytes:
    """sb_dict: {kategori: df_scoreboard} -> satu PDF landscape berisi semua kategori."""
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=landscape(A4), topMargin=1 * cm, bottomMargin=1 * cm,
                             leftMargin=1 * cm, rightMargin=1 * cm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("title", parent=styles["Title"], fontSize=14, textColor=rl_colors.HexColor("#0f766e"))
    elements = []
    for kategori, df_sb in sb_dict.items():
        elements.append(Paragraph(_SCOREBOARD_KATEGORI_LABEL.get(kategori, kategori), title_style))
        elements.append(Spacer(1, 8))
        if df_sb.empty:
            elements.append(Paragraph("Tidak ada data untuk kategori ini.", styles["Normal"]))
        else:
            col_labels = [_SCOREBOARD_GROUPS.get(c, c) for c in _SCOREBOARD_COL_ORDER]
            data = [col_labels]
            for _, r in df_sb.iterrows():
                data.append([_fmt_scoreboard_cell(c, r.get(c)) for c in _SCOREBOARD_COL_ORDER])
            table = Table(data, repeatRows=1)
            style_cmds = [
                ("BACKGROUND", (0, 0), (-1, 0), rl_colors.HexColor("#0f766e")),
                ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.4, rl_colors.HexColor("#d1d5db")),
                ("ALIGN", (1, 0), (-1, -1), "CENTER"),
                ("FONTSIZE", (0, 0), (-1, -1), 6.5),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ]
            for i, (_, r) in enumerate(df_sb.iterrows(), start=1):
                if str(r["Cabang"]).upper() == "TOTAL":
                    style_cmds.append(("BACKGROUND", (0, i), (-1, i), rl_colors.HexColor("#f0fdfa")))
                    style_cmds.append(("FONTNAME", (0, i), (-1, i), "Helvetica-Bold"))
            table.setStyle(TableStyle(style_cmds))
            elements.append(table)
        elements.append(Spacer(1, 18))
    doc.build(elements)
    buf.seek(0)
    return buf.read()


# ========================= Progress Ring / Charts ==========================

def render_progress_ring(pct: float, label: str, size: int = 130) -> str:
    pct_clamped = max(0.0, min(1.0, pct if pct is not None else 0.0))
    color = pencapaian_color(pct)
    deg = pct_clamped * 360
    pct_text = format_percent(pct) if pct is not None else "-"
    return f"""
    <div style="display:flex;flex-direction:column;align-items:center;margin:6px;">
        <div style="width:{size}px;height:{size}px;border-radius:50%;
            background:conic-gradient({color} {deg}deg, #e5e7eb {deg}deg);
            display:flex;align-items:center;justify-content:center;">
            <div style="width:{size-24}px;height:{size-24}px;border-radius:50%;background:white;
                display:flex;align-items:center;justify-content:center;flex-direction:column;">
                <span style="font-size:1.1em;font-weight:800;color:{color};">{pct_text}</span>
            </div>
        </div>
        <span style="margin-top:6px;font-size:0.85em;font-weight:600;color:#374151;text-align:center;">{label}</span>
    </div>
    """


def render_contribution_pie(labels: list, values: list, colors: list, title: str = ""):
    fig = go.Figure(data=[go.Pie(labels=labels, values=values, marker=dict(colors=colors), hole=0.45)])
    fig.update_layout(title=title, showlegend=True, height=320, margin=dict(t=40, b=10, l=10, r=10))
    return fig


def build_daily_progress(df_main: pd.DataFrame, tanggal_acuan: date, selected_branches=None) -> pd.DataFrame:
    if df_main is None or df_main.empty or tanggal_acuan is None:
        return pd.DataFrame(columns=["Tanggal", "Omset"])
    start, end, *_ = _quarter_bounds(tanggal_acuan)
    d = df_main[(df_main["Tanggal"] >= start) & (df_main["Tanggal"] <= tanggal_acuan)]
    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]
    if d.empty:
        return pd.DataFrame(columns=["Tanggal", "Omset"])
    g = d.groupby("Tanggal")["Omset"].sum().reset_index().sort_values("Tanggal")
    return g


def render_daily_progress_chart(g: pd.DataFrame):
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=g["Tanggal"], y=g["Omset"], mode="lines+markers", line=dict(color="#0f766e", width=2),
                              marker=dict(size=5), fill="tozeroy", fillcolor="rgba(15,118,110,0.1)"))
    fig.update_layout(height=320, margin=dict(t=20, b=10, l=10, r=10), xaxis_title="Tanggal", yaxis_title="Omset")
    return fig


def build_daily_history(df_main: pd.DataFrame, selected_branches=None, n_days: int = 30) -> pd.DataFrame:
    if df_main is None or df_main.empty:
        return pd.DataFrame(columns=["Tanggal", "Omset"])
    d = df_main.copy()
    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]
    if d.empty:
        return pd.DataFrame(columns=["Tanggal", "Omset"])
    g = d.groupby("Tanggal")["Omset"].sum().reset_index().sort_values("Tanggal")
    return g.tail(n_days)


def render_daily_history_chart(g: pd.DataFrame):
    fig = go.Figure()
    fig.add_trace(go.Bar(x=g["Tanggal"], y=g["Omset"], marker_color="#0f766e"))
    fig.update_layout(height=300, margin=dict(t=20, b=10, l=10, r=10), xaxis_title="Tanggal", yaxis_title="Omset")
    return fig


# ========================= 6 Pilar ==========================

def build_pilar_summary(df_main: pd.DataFrame, tanggal_acuan: date, selected_branches=None) -> pd.DataFrame:
    cols = ["Pilar", "Omset", "Qty"]
    if df_main is None or df_main.empty or tanggal_acuan is None or "Pilar" not in df_main.columns:
        return pd.DataFrame(columns=cols)
    start, end, *_ = _quarter_bounds(tanggal_acuan)
    d = df_main[(df_main["Tanggal"] >= start) & (df_main["Tanggal"] <= tanggal_acuan)]
    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]
    if d.empty:
        return pd.DataFrame(columns=cols)
    g = d.groupby("Pilar").agg(Omset=("Omset", "sum"), Qty=("Qty", "sum")).reset_index()
    g["_order"] = g["Pilar"].apply(lambda p: PILAR_ORDER.index(p) if p in PILAR_ORDER else 99)
    g = g.sort_values("_order").drop(columns="_order").reset_index(drop=True)
    return g


def build_pilar_by_branch(df_main: pd.DataFrame, tanggal_acuan: date, selected_branches=None) -> pd.DataFrame:
    cols = ["Cabang", "Pilar", "Omset", "Qty"]
    if df_main is None or df_main.empty or tanggal_acuan is None or "Pilar" not in df_main.columns:
        return pd.DataFrame(columns=cols)
    start, end, *_ = _quarter_bounds(tanggal_acuan)
    d = df_main[(df_main["Tanggal"] >= start) & (df_main["Tanggal"] <= tanggal_acuan)]
    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]
    if d.empty:
        return pd.DataFrame(columns=cols)
    g = d.groupby(["Cabang", "Pilar"]).agg(Omset=("Omset", "sum"), Qty=("Qty", "sum")).reset_index()
    return g


def generate_pilar_insights(pilar_summary: pd.DataFrame) -> list:
    insights = []
    if pilar_summary.empty:
        return insights
    total = pilar_summary["Omset"].sum()
    if total <= 0:
        return insights
    for _, r in pilar_summary.iterrows():
        share = r["Omset"] / total
        if share < 0.03:
            insights.append({
                "level": "warn",
                "category": r["Pilar"],
                "title": f"{_pilar_label(r['Pilar'])}: kontribusi masih sangat kecil ({format_percent(share)})",
                "problem": f"Omset pilar ini hanya {format_rupiah(r['Omset'])} dari total {format_rupiah(total)}.",
                "online": ["Promosikan pilar ini lebih intensif di sosial media & katalog online"],
                "offline": ["Latih tim sales untuk cross-selling pilar ini ke pelanggan yang datang"],
            })
    return insights


def render_pilar_kpi_card(pilar: str, omset: float, qty: float) -> str:
    icon = PILAR_ICONS.get(pilar, "🧩")
    color = PILAR_COLORS.get(pilar, "#0f766e")
    show_qty = pilar in _PILAR_SHOW_QTY
    qty_html = f'<div style="font-size:0.8em;color:#6b7280;">Qty: {format_number(qty)}</div>' if show_qty else ""
    return f"""
    <div style="border:2px solid {color};border-radius:10px;padding:10px 14px;text-align:center;background:white;">
        <div style="font-size:1.6em;">{icon}</div>
        <div style="font-size:0.82em;font-weight:700;color:{color};margin:4px 0;">{_pilar_label(pilar)}</div>
        <div style="font-size:1.05em;font-weight:800;color:#111827;">{format_rupiah(omset)}</div>
        {qty_html}
    </div>
    """


def render_pilar_table_html(g: pd.DataFrame) -> str:
    if g.empty:
        return "<i>Tidak ada data 6 Pilar untuk periode ini.</i>"
    rows_html = ""
    for _, r in g.iterrows():
        rows_html += f"""<tr>
            <td style="padding:6px 10px;border:1px solid #d1d5db;">{PILAR_ICONS.get(r['Pilar'],'')} {_pilar_label(r['Pilar'])}</td>
            <td style="padding:6px 10px;border:1px solid #d1d5db;text-align:right;">{format_rupiah(r['Omset'])}</td>
            <td style="padding:6px 10px;border:1px solid #d1d5db;text-align:right;">{format_number(r['Qty']) if r['Pilar'] in _PILAR_SHOW_QTY else '-'}</td>
        </tr>"""
    return f"""
    <table style="border-collapse:collapse;width:100%;font-size:0.9em;">
        <thead><tr style="background:#f0fdfa;">
            <th style="padding:6px 10px;border:1px solid #d1d5db;text-align:left;">Pilar</th>
            <th style="padding:6px 10px;border:1px solid #d1d5db;">Omset</th>
            <th style="padding:6px 10px;border:1px solid #d1d5db;">Qty</th>
        </tr></thead>
        <tbody>{rows_html}</tbody>
    </table>
    """


def render_pilar_summary_table_html(g_branch: pd.DataFrame) -> str:
    if g_branch.empty:
        return "<i>Tidak ada data.</i>"
    pivot = g_branch.pivot_table(index="Cabang", columns="Pilar", values="Omset", aggfunc="sum", fill_value=0)
    pivot = pivot.reindex(columns=[p for p in PILAR_ORDER if p in pivot.columns])
    header = "".join(f'<th style="padding:5px 8px;border:1px solid #d1d5db;">{_pilar_label(p)}</th>' for p in pivot.columns)
    rows_html = ""
    for cabang in order_branches(list(pivot.index)):
        if cabang not in pivot.index:
            continue
        row = pivot.loc[cabang]
        cells = "".join(f'<td style="padding:5px 8px;border:1px solid #d1d5db;text-align:right;">{format_rupiah(v)}</td>' for v in row)
        rows_html += f'<tr><td style="padding:5px 8px;border:1px solid #d1d5db;font-weight:600;">{cabang}</td>{cells}</tr>'
    return f"""
    <table style="border-collapse:collapse;width:100%;font-size:0.82em;">
        <thead><tr style="background:#f0fdfa;"><th style="padding:5px 8px;border:1px solid #d1d5db;text-align:left;">Cabang</th>{header}</tr></thead>
        <tbody>{rows_html}</tbody>
    </table>
    """


# ========================= Kontribusi Marketing Corporate vs Sales Retail ==========================

_MC_CATEGORY_ORDER = ["Marketing Corporate", "Sales Retail"]
_MC_CATEGORY_LABELS = {"Marketing Corporate": "Marketing Corporate", "Sales Retail": "Sales Retail"}
_MC_CATEGORY_ICONS = {"Marketing Corporate": "🤝", "Sales Retail": "🏪"}
_MC_CATEGORY_COLORS = {"Marketing Corporate": "#7c3aed", "Sales Retail": "#0f766e"}


def _mc_filter_bulan_berjalan(df_main: pd.DataFrame, tanggal_acuan: date) -> pd.DataFrame:
    if df_main is None or df_main.empty or tanggal_acuan is None:
        return pd.DataFrame()
    start = date(tanggal_acuan.year, tanggal_acuan.month, 1)
    return df_main[(df_main["Tanggal"] >= start) & (df_main["Tanggal"] <= tanggal_acuan)]


def build_mc_contribution_summary(df_main: pd.DataFrame, tanggal_acuan: date, selected_branches=None) -> pd.DataFrame:
    cols = ["Kelompok", "Omset"]
    if df_main is None or df_main.empty or "PenjualKelompok" not in df_main.columns:
        return pd.DataFrame(columns=cols)
    start, end, *_ = _quarter_bounds(tanggal_acuan)
    d = df_main[(df_main["Tanggal"] >= start) & (df_main["Tanggal"] <= tanggal_acuan)]
    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]
    if d.empty:
        return pd.DataFrame(columns=cols)
    g = d.groupby("PenjualKelompok")["Omset"].sum().reset_index()
    g.columns = ["Kelompok", "Omset"]
    return g


def build_mc_person_table(df_main: pd.DataFrame, tanggal_acuan: date, selected_branches=None) -> pd.DataFrame:
    cols = ["NamaPenjual", "Cabang", "Omset"]
    if df_main is None or df_main.empty or "PenjualKelompok" not in df_main.columns:
        return pd.DataFrame(columns=cols)
    start, end, *_ = _quarter_bounds(tanggal_acuan)
    d = df_main[(df_main["Tanggal"] >= start) & (df_main["Tanggal"] <= tanggal_acuan) & (df_main["PenjualKelompok"] == "Marketing Corporate")]
    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]
    if d.empty or "NamaPenjual" not in d.columns:
        return pd.DataFrame(columns=cols)
    g = d.groupby(["NamaPenjual", "Cabang"])["Omset"].sum().reset_index()
    return g.sort_values("Omset", ascending=False)


def build_retail_by_branch(df_main: pd.DataFrame, tanggal_acuan: date, selected_branches=None) -> pd.DataFrame:
    cols = ["Cabang", "Omset"]
    if df_main is None or df_main.empty or "PenjualKelompok" not in df_main.columns:
        return pd.DataFrame(columns=cols)
    start, end, *_ = _quarter_bounds(tanggal_acuan)
    d = df_main[(df_main["Tanggal"] >= start) & (df_main["Tanggal"] <= tanggal_acuan) & (df_main["PenjualKelompok"] == "Sales Retail")]
    if selected_branches:
        d = d[d["Cabang"].isin(selected_branches)]
    if d.empty:
        return pd.DataFrame(columns=cols)
    g = d.groupby("Cabang")["Omset"].sum().reset_index()
    return g


def render_mc_contribution_card(kelompok: str, omset: float, total: float) -> str:
    icon = _MC_CATEGORY_ICONS.get(kelompok, "📊")
    color = _MC_CATEGORY_COLORS.get(kelompok, "#374151")
    share = (omset / total) if total else 0.0
    return f"""
    <div style="border:2px solid {color};border-radius:10px;padding:12px 16px;text-align:center;background:white;">
        <div style="font-size:1.8em;">{icon}</div>
        <div style="font-size:0.9em;font-weight:700;color:{color};margin:4px 0;">{_MC_CATEGORY_LABELS.get(kelompok, kelompok)}</div>
        <div style="font-size:1.2em;font-weight:800;color:#111827;">{format_rupiah(omset)}</div>
        <div style="font-size:0.85em;color:#6b7280;">{format_percent(share)} dari total</div>
    </div>
    """


def render_mc_split_donut(g: pd.DataFrame):
    if g.empty:
        return None
    labels = [_MC_CATEGORY_LABELS.get(k, k) for k in g["Kelompok"]]
    colors_list = [_MC_CATEGORY_COLORS.get(k, "#9ca3af") for k in g["Kelompok"]]
    fig = go.Figure(data=[go.Pie(labels=labels, values=g["Omset"], marker=dict(colors=colors_list), hole=0.5)])
    fig.update_layout(height=300, margin=dict(t=20, b=10, l=10, r=10))
    return fig


def render_mc_person_table_html(g: pd.DataFrame) -> str:
    if g.empty:
        return "<i>Tidak ada data Marketing Corporate untuk periode ini.</i>"
    rows_html = ""
    for _, r in g.iterrows():
        rows_html += f"""<tr>
            <td style="padding:6px 10px;border:1px solid #d1d5db;">{r['NamaPenjual']}</td>
            <td style="padding:6px 10px;border:1px solid #d1d5db;">{r['Cabang']}</td>
            <td style="padding:6px 10px;border:1px solid #d1d5db;text-align:right;">{format_rupiah(r['Omset'])}</td>
        </tr>"""
    return f"""
    <table style="border-collapse:collapse;width:100%;font-size:0.9em;">
        <thead><tr style="background:#f5f3ff;">
            <th style="padding:6px 10px;border:1px solid #d1d5db;text-align:left;">Nama Marketing</th>
            <th style="padding:6px 10px;border:1px solid #d1d5db;text-align:left;">Cabang</th>
            <th style="padding:6px 10px;border:1px solid #d1d5db;">Omset</th>
        </tr></thead>
        <tbody>{rows_html}</tbody>
    </table>
    """


def render_retail_by_branch_table_html(g: pd.DataFrame) -> str:
    if g.empty:
        return "<i>Tidak ada data Sales Retail untuk periode ini.</i>"
    rows_html = ""
    for _, r in g.iterrows():
        rows_html += f"""<tr>
            <td style="padding:6px 10px;border:1px solid #d1d5db;font-weight:600;">{r['Cabang']}</td>
            <td style="padding:6px 10px;border:1px solid #d1d5db;text-align:right;">{format_rupiah(r['Omset'])}</td>
        </tr>"""
    return f"""
    <table style="border-collapse:collapse;width:100%;font-size:0.9em;">
        <thead><tr style="background:#f0fdfa;">
            <th style="padding:6px 10px;border:1px solid #d1d5db;text-align:left;">Cabang</th>
            <th style="padding:6px 10px;border:1px solid #d1d5db;">Omset Retail</th>
        </tr></thead>
        <tbody>{rows_html}</tbody>
    </table>
    """


_MC_MARKETING_ACTION_PLAN = {
    "online": ["Follow-up leads corporate via WhatsApp Business/LinkedIn", "Email penawaran kerja sama ke database corporate lama"],
    "offline": ["Kunjungan langsung ke kantor/instansi calon klien corporate", "Ikuti pameran/expo B2B di area cabang"],
}


def generate_mc_insights(g_summary: pd.DataFrame) -> list:
    insights = []
    if g_summary.empty:
        return insights
    total = g_summary["Omset"].sum()
    if total <= 0:
        return insights
    mc_row = g_summary[g_summary["Kelompok"] == "Marketing Corporate"]
    if not mc_row.empty:
        share = mc_row.iloc[0]["Omset"] / total
        if share < 0.1:
            insights.append({
                "level": "warn", "category": "Marketing Corporate",
                "title": f"Kontribusi Marketing Corporate masih rendah ({format_percent(share)})",
                "problem": f"Omset Corporate {format_rupiah(mc_row.iloc[0]['Omset'])} dari total {format_rupiah(total)}.",
                "online": _MC_MARKETING_ACTION_PLAN["online"], "offline": _MC_MARKETING_ACTION_PLAN["offline"],
            })
    return insights


def generate_retail_branch_insights(g_retail: pd.DataFrame) -> list:
    insights = []
    if g_retail.empty:
        return insights
    avg = g_retail["Omset"].mean()
    if avg <= 0:
        return insights
    for _, r in g_retail.iterrows():
        if r["Omset"] < 0.7 * avg:
            insights.append({
                "level": "bad", "category": r["Cabang"],
                "title": f"{r['Cabang']}: Omset Retail jauh di bawah rata-rata cabang",
                "problem": f"Omset {format_rupiah(r['Omset'])}, rata-rata cabang lain {format_rupiah(avg)}.",
                "online": ["Tingkatkan promosi lokal & konten media sosial cabang"],
                "offline": ["Evaluasi layanan pelanggan & display produk toko"],
            })
    return insights


# ========================= Ledger (Riwayat Upload) ==========================

_HISTORY_LOG_COLUMNS = ["Tanggal", "Cabang", "Kategori", "Omset", "Timestamp"]
_LOG_PATH = os.path.join(DATA_DIR, "log", "upload_log.csv")


def _read_log() -> pd.DataFrame:
    if not os.path.exists(_LOG_PATH):
        return pd.DataFrame(columns=_HISTORY_LOG_COLUMNS)
    try:
        df = pd.read_csv(_LOG_PATH)
        if "Tanggal" in df.columns:
            df["Tanggal"] = pd.to_datetime(df["Tanggal"]).dt.date
        return df
    except Exception:
        return pd.DataFrame(columns=_HISTORY_LOG_COLUMNS)


def _upsert_log(df_main: pd.DataFrame):
    if df_main is None or df_main.empty:
        return
    os.makedirs(os.path.dirname(_LOG_PATH), exist_ok=True)
    g = df_main.groupby(["Tanggal", "Cabang", "Kategori"])["Omset"].sum().reset_index()
    g["Timestamp"] = datetime.now().isoformat()
    g.to_csv(_LOG_PATH, index=False)
    try:
        if _GH_ENABLED:
            github_upload_file(f"data/log/{os.path.basename(_LOG_PATH)}", open(_LOG_PATH, "rb").read())
    except Exception:
        pass


def build_upload_log(df_main: pd.DataFrame) -> pd.DataFrame:
    _upsert_log(df_main)
    return _read_log()


def build_corp_upload_log(df_corp: pd.DataFrame) -> pd.DataFrame:
    """Stub: log riwayat upload corporate (belum ada kebutuhan spesifik)."""
    return pd.DataFrame()


def compute_corp_hari_ini(df_corp: pd.DataFrame, tanggal_acuan: date) -> float:
    """Stub fallback: pakai total omset corporate manual dibagi hari berjalan bulan ini."""
    if df_corp is None or df_corp.empty or "Omset Corporate" not in df_corp.columns:
        return 0.0
    total = float(df_corp["Omset Corporate"].sum())
    if tanggal_acuan is None:
        return total
    days_in_month = tanggal_acuan.day
    return total / days_in_month if days_in_month else total


# ========================= PPTX / PDF Full Report Export ==========================

def _pptx_add_title_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def _pptx_add_table_slide(prs, title: str, df: pd.DataFrame, fmt_fn=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if df is None or df.empty:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        tx.text_frame.text = "Tidak ada data."
        return slide
    rows, cols = df.shape[0] + 1, df.shape[1]
    left, top, width, height = Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.35 * (rows + 1))
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = str(col)
    for i, (_, r) in enumerate(df.iterrows(), start=1):
        for j, col in enumerate(df.columns):
            val = r[col]
            txt = fmt_fn(col, val) if fmt_fn else str(val)
            table.cell(i, j).text = txt
    return slide


def _build_report_sections(df_main, sb_dict, walkin_df, pilar_summary, mc_summary):
    sections = []
    sections.append(("Scoreboard Omset All", sb_dict.get("Omset All", pd.DataFrame())))
    sections.append(("Scoreboard Service", sb_dict.get("Service", pd.DataFrame())))
    sections.append(("Scoreboard Gadget & Aksesoris", sb_dict.get("Gadget & Aksesoris", pd.DataFrame())))
    sections.append(("Walk-in per Cabang", walkin_df))
    sections.append(("6 Pilar MFlash", pilar_summary))
    sections.append(("Kontribusi Marketing Corporate vs Retail", mc_summary))
    return sections


def generate_pptx_report(df_main, sb_dict, walkin_df, pilar_summary, mc_summary, periode_label=""):
    prs = Presentation()
    _pptx_add_title_slide(prs, "Laporan Dashboard Omset MFlash", periode_label)
    sections = _build_report_sections(df_main, sb_dict, walkin_df, pilar_summary, mc_summary)
    for title, df in sections:
        _pptx_add_table_slide(prs, title, df, fmt_fn=lambda c, v: _fmt_scoreboard_cell(c, v) if c in _SCOREBOARD_COL_ORDER else (
            format_rupiah(v) if isinstance(v, (int, float)) else str(v)))
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def generate_pdf_report(df_main, sb_dict, walkin_df, pilar_summary, mc_summary, periode_label=""):
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=landscape(A4), topMargin=1 * cm, bottomMargin=1 * cm,
                             leftMargin=1 * cm, rightMargin=1 * cm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("title", parent=styles["Title"], fontSize=16, textColor=rl_colors.HexColor("#0f766e"))
    sub_style = ParagraphStyle("sub", parent=styles["Normal"], fontSize=10, textColor=rl_colors.HexColor("#374151"))
    elements = [Paragraph("Laporan Dashboard Omset MFlash", title_style), Paragraph(periode_label, sub_style), Spacer(1, 14)]

    sections = _build_report_sections(df_main, sb_dict, walkin_df, pilar_summary, mc_summary)
    for title, df in sections:
        elements.append(Paragraph(title, ParagraphStyle("h2", parent=styles["Heading2"], textColor=rl_colors.HexColor("#0f766e"))))
        elements.append(Spacer(1, 6))
        if df is None or df.empty:
            elements.append(Paragraph("Tidak ada data.", styles["Normal"]))
        else:
            data = [list(df.columns)]
            for _, r in df.iterrows():
                row = []
                for c in df.columns:
                    v = r[c]
                    if c in _SCOREBOARD_COL_ORDER:
                        row.append(_fmt_scoreboard_cell(c, v))
                    elif isinstance(v, (int, float)):
                        row.append(format_rupiah(v) if abs(v) > 1000 else format_number(v))
                    else:
                        row.append(str(v))
                data.append(row)
            table = Table(data, repeatRows=1)
            table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), rl_colors.HexColor("#0f766e")),
                ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.4, rl_colors.HexColor("#d1d5db")),
                ("FONTSIZE", (0, 0), (-1, -1), 6.5),
                ("ALIGN", (1, 0), (-1, -1), "CENTER"),
            ]))
            elements.append(table)
        elements.append(Spacer(1, 16))
    doc.build(elements)
    buf.seek(0)
    return buf.read()


# ========================= UI ==========================

for _d in [MAIN_DATA_DIR, ADS_DATA_DIR, WALKIN_DATA_DIR, TARGET_DATA_DIR, CORP_DATA_DIR, LOG_DIR]:
    os.makedirs(_d, exist_ok=True)

if _GH_ENABLED and not st.session_state.get("_gh_synced"):
    try:
        sync_data_from_github()
    except Exception:
        pass
    st.session_state["_gh_synced"] = True

_dedupe_main_files()

st.markdown(
    f"""
    <div style="display:flex;align-items:center;gap:16px;margin-bottom:10px;">
        <img src="data:image/png;base64,{LOGO_BASE64}" style="height:64px;" />
        <div>
            <div style="font-size:1.5em;font-weight:800;color:#0f766e;">Dashboard Omset MFlash</div>
            <div style="font-size:0.9em;color:#6b7280;">Monitoring Omset, Iklan, Walk-in, 6 Pilar &amp; Kontribusi Marketing</div>
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

with st.sidebar:
    st.header("📁 Upload Data")

    with st.expander("📊 Data Omset (Faktur Penjualan)", expanded=False):
        up_main = st.file_uploader("Upload file Excel Omset", type=["xlsx", "xls"], accept_multiple_files=True, key="up_main")
        if up_main:
            for f in up_main:
                dest = os.path.join(MAIN_DATA_DIR, f.name)
                with open(dest, "wb") as out:
                    out.write(f.getbuffer())
                if _GH_ENABLED:
                    try:
                        github_upload_file(f"data/main/{f.name}", f.getbuffer())
                    except Exception:
                        pass
            _dedupe_main_files()
            st.success(f"{len(up_main)} file diunggah.")
        existing_main = sorted(os.listdir(MAIN_DATA_DIR)) if os.path.isdir(MAIN_DATA_DIR) else []
        for fn in existing_main:
            c1, c2 = st.columns([4, 1])
            c1.caption(fn)
            if c2.button("🗑️", key=f"del_main_{fn}"):
                os.remove(os.path.join(MAIN_DATA_DIR, fn))
                if _GH_ENABLED:
                    try:
                        github_delete_file(f"data/main/{fn}")
                    except Exception:
                        pass
                st.rerun()

    with st.expander("📢 Data Iklan (Meta Ads)", expanded=False):
        up_ads = st.file_uploader("Upload file Excel Ads", type=["xlsx", "xls", "csv"], accept_multiple_files=True, key="up_ads")
        if up_ads:
            for f in up_ads:
                dest = os.path.join(ADS_DATA_DIR, f.name)
                with open(dest, "wb") as out:
                    out.write(f.getbuffer())
                if _GH_ENABLED:
                    try:
                        github_upload_file(f"data/ads/{f.name}", f.getbuffer())
                    except Exception:
                        pass
            st.success(f"{len(up_ads)} file diunggah.")
        existing_ads = sorted(os.listdir(ADS_DATA_DIR)) if os.path.isdir(ADS_DATA_DIR) else []
        for fn in existing_ads:
            c1, c2 = st.columns([4, 1])
            c1.caption(fn)
            if c2.button("🗑️", key=f"del_ads_{fn}"):
                os.remove(os.path.join(ADS_DATA_DIR, fn))
                if _GH_ENABLED:
                    try:
                        github_delete_file(f"data/ads/{fn}")
                    except Exception:
                        pass
                st.rerun()

    with st.expander("🚶 Data Walk-in", expanded=False):
        up_walkin = st.file_uploader("Upload file Excel Walk-in", type=["xlsx", "xls"], accept_multiple_files=True, key="up_walkin")
        if up_walkin:
            for f in up_walkin:
                dest = os.path.join(WALKIN_DATA_DIR, f.name)
                with open(dest, "wb") as out:
                    out.write(f.getbuffer())
                if _GH_ENABLED:
                    try:
                        github_upload_file(f"data/walkin/{f.name}", f.getbuffer())
                    except Exception:
                        pass
            st.success(f"{len(up_walkin)} file diunggah.")
        existing_walkin = sorted(os.listdir(WALKIN_DATA_DIR)) if os.path.isdir(WALKIN_DATA_DIR) else []
        for fn in existing_walkin:
            c1, c2 = st.columns([4, 1])
            c1.caption(fn)
            if c2.button("🗑️", key=f"del_walkin_{fn}"):
                os.remove(os.path.join(WALKIN_DATA_DIR, fn))
                if _GH_ENABLED:
                    try:
                        github_delete_file(f"data/walkin/{fn}")
                    except Exception:
                        pass
                st.rerun()

    with st.expander("🎯 Target Omset (opsional)", expanded=False):
        st.download_button("⬇️ Download Template Target", data=make_target_template(),
                            file_name="template_target.xlsx", key="dl_target_tpl")
        up_target = st.file_uploader("Upload file Target", type=["xlsx", "xls"], accept_multiple_files=True, key="up_target")
        if up_target:
            for f in up_target:
                dest = os.path.join(TARGET_DATA_DIR, f.name)
                with open(dest, "wb") as out:
                    out.write(f.getbuffer())
                if _GH_ENABLED:
                    try:
                        github_upload_file(f"data/target/{f.name}", f.getbuffer())
                    except Exception:
                        pass
            st.success(f"{len(up_target)} file diunggah.")
        existing_target = sorted(os.listdir(TARGET_DATA_DIR)) if os.path.isdir(TARGET_DATA_DIR) else []
        for fn in existing_target:
            c1, c2 = st.columns([4, 1])
            c1.caption(fn)
            if c2.button("🗑️", key=f"del_target_{fn}"):
                os.remove(os.path.join(TARGET_DATA_DIR, fn))
                if _GH_ENABLED:
                    try:
                        github_delete_file(f"data/target/{fn}")
                    except Exception:
                        pass
                st.rerun()

    with st.expander("🤝 Data Corporate (opsional)", expanded=False):
        st.download_button("⬇️ Download Template Corporate", data=make_corporate_template(),
                            file_name="template_corporate.xlsx", key="dl_corp_tpl")
        up_corp = st.file_uploader("Upload file Corporate", type=["xlsx", "xls"], accept_multiple_files=True, key="up_corp")
        if up_corp:
            for f in up_corp:
                dest = os.path.join(CORP_DATA_DIR, f.name)
                with open(dest, "wb") as out:
                    out.write(f.getbuffer())
                if _GH_ENABLED:
                    try:
                        github_upload_file(f"data/corp/{f.name}", f.getbuffer())
                    except Exception:
                        pass
            st.success(f"{len(up_corp)} file diunggah.")
        existing_corp = sorted(os.listdir(CORP_DATA_DIR)) if os.path.isdir(CORP_DATA_DIR) else []
        for fn in existing_corp:
            c1, c2 = st.columns([4, 1])
            c1.caption(fn)
            if c2.button("🗑️", key=f"del_corp_{fn}"):
                os.remove(os.path.join(CORP_DATA_DIR, fn))
                if _GH_ENABLED:
                    try:
                        github_delete_file(f"data/corp/{fn}")
                    except Exception:
                        pass
                st.rerun()

df_main = load_all_main_data()
df_ads = load_all_ads_data()
df_walkin = load_all_walkin_data()

target_map = {k: {} for k in SCOREBOARD_KATEGORI}
target_files = sorted(os.listdir(TARGET_DATA_DIR)) if os.path.isdir(TARGET_DATA_DIR) else []
if target_files:
    for fn in target_files:
        df_t = load_target_data(os.path.join(TARGET_DATA_DIR, fn))
        if df_t.empty:
            continue
        cols_lower = {c.lower(): c for c in df_t.columns}
        kat_col = cols_lower.get("kategori")
        cab_col = cols_lower.get("cabang")
        val_col = None
        for c in df_t.columns:
            if "target" in c.lower() and "omset" in c.lower():
                val_col = c
                break
        if kat_col and cab_col and val_col:
            for _, r in df_t.iterrows():
                kat = str(r[kat_col]).strip()
                cab = str(r[cab_col]).strip().upper()
                if kat in target_map:
                    try:
                        target_map[kat][cab] = float(r[val_col])
                    except (ValueError, TypeError):
                        pass
else:
    main_files = sorted(os.listdir(MAIN_DATA_DIR)) if os.path.isdir(MAIN_DATA_DIR) else []
    for fn in main_files:
        fpath = os.path.join(MAIN_DATA_DIR, fn)
        extracted = extract_scoreboard_target(fpath)
        for kat, mapping in extracted.items():
            for cab, val in mapping.items():
                if cab not in target_map[kat]:
                    target_map[kat][cab] = val

df_corp = pd.DataFrame()
corp_files = sorted(os.listdir(CORP_DATA_DIR)) if os.path.isdir(CORP_DATA_DIR) else []
if corp_files:
    frames = [load_corporate_data(os.path.join(CORP_DATA_DIR, fn)) for fn in corp_files]
    frames = [f for f in frames if not f.empty]
    if frames:
        df_corp = pd.concat(frames, ignore_index=True)

if df_main.empty:
    st.info("👋 Silakan unggah minimal 1 file Data Omset (Faktur Penjualan) di sidebar untuk mulai menggunakan dashboard.")

default_date = df_main["Tanggal"].max() if not df_main.empty and "Tanggal" in df_main.columns else date.today()
if pd.isna(default_date):
    default_date = date.today()
if isinstance(default_date, pd.Timestamp):
    default_date = default_date.date()

col_f1, col_f2 = st.columns([1, 2])
with col_f1:
    tanggal_acuan = st.date_input("📅 Tanggal Acuan", value=default_date, key="tanggal_acuan")
with col_f2:
    all_branches_present = order_branches(df_main["Cabang"].unique()) if not df_main.empty else BRANCH_ORDER
    selected_branches = st.multiselect("🏢 Filter Cabang", options=all_branches_present, default=all_branches_present, key="selected_branches")

q_start, q_end, q_total_hari, q_hari_berjalan, q_sisa_hari = _quarter_bounds(tanggal_acuan)
periode_label = f"{BULAN_ID.get(tanggal_acuan.month,'')} {tanggal_acuan.year}"
quarter_period_label = f"{q_start.day} {BULAN_ID.get(q_start.month,'')[:3]} - {tanggal_acuan.day} {BULAN_ID.get(tanggal_acuan.month,'')[:3]} {tanggal_acuan.year}"

sb_dict = {}
for kat in SCOREBOARD_KATEGORI:
    sb_raw = build_scoreboard(df_main, tanggal_acuan, target_map, kat, selected_branches)
    sb_dict[kat] = _finalize_scoreboard(sb_raw)

walkin_agg_current = aggregate_walkin_current_period(df_walkin, tanggal_acuan, selected_branches) if not df_walkin.empty else pd.DataFrame()
walkin_agg_monthly = aggregate_walkin_monthly(df_walkin) if not df_walkin.empty else pd.DataFrame()

pilar_summary = build_pilar_summary(df_main, tanggal_acuan, selected_branches)
pilar_by_branch = build_pilar_by_branch(df_main, tanggal_acuan, selected_branches)

mc_summary = build_mc_contribution_summary(df_main, tanggal_acuan, selected_branches)
mc_person_table = build_mc_person_table(df_main, tanggal_acuan, selected_branches)
retail_by_branch = build_retail_by_branch(df_main, tanggal_acuan, selected_branches)

if not df_main.empty:
    build_upload_log(df_main)

tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "🏠 Ringkasan", "🏆 Scoreboard", "📢 Iklan", "🚶 Walk-in", "🧩 6 Pilar", "🤝 Kontribusi MC",
])

with tab1:
    st.subheader(f"🏠 Ringkasan — {periode_label}")
    if df_main.empty:
        st.warning("Belum ada data Omset.")
    else:
        omset_all_row = sb_dict["Omset All"][sb_dict["Omset All"]["Cabang"] == "TOTAL"]
        service_row = sb_dict["Service"][sb_dict["Service"]["Cabang"] == "TOTAL"]
        gadget_row = sb_dict["Gadget & Aksesoris"][sb_dict["Gadget & Aksesoris"]["Cabang"] == "TOTAL"]

        c1, c2, c3 = st.columns(3)
        with c1:
            v = omset_all_row.iloc[0] if not omset_all_row.empty else None
            st.markdown(render_kpi_card("S/D Hari Ini (Omset All)", format_rupiah(v["SdHariIni"]) if v is not None else "Rp 0", "#0f766e", "💰"), unsafe_allow_html=True)
        with c2:
            v = service_row.iloc[0] if not service_row.empty else None
            st.markdown(render_kpi_card("S/D Hari Ini (Service)", format_rupiah(v["SdHariIni"]) if v is not None else "Rp 0", "#2563eb", "🔧"), unsafe_allow_html=True)
        with c3:
            v = gadget_row.iloc[0] if not gadget_row.empty else None
            st.markdown(render_kpi_card("S/D Hari Ini (Gadget & Aksesoris)", format_rupiah(v["SdHariIni"]) if v is not None else "Rp 0", "#d97706", "📱"), unsafe_allow_html=True)

        st.markdown("###### % Pencapaian Target")
        r1, r2, r3 = st.columns(3)
        with r1:
            v = omset_all_row.iloc[0] if not omset_all_row.empty else None
            st.markdown(render_progress_ring(v["PctPencapaian"] if v is not None else None, "Omset All"), unsafe_allow_html=True)
        with r2:
            v = service_row.iloc[0] if not service_row.empty else None
            st.markdown(render_progress_ring(v["PctPencapaian"] if v is not None else None, "Service"), unsafe_allow_html=True)
        with r3:
            v = gadget_row.iloc[0] if not gadget_row.empty else None
            st.markdown(render_progress_ring(v["PctPencapaian"] if v is not None else None, "Gadget & Aksesoris"), unsafe_allow_html=True)

        st.markdown("###### Progres Harian (Kuartal Berjalan)")
        g_daily = build_daily_progress(df_main, tanggal_acuan, selected_branches)
        if not g_daily.empty:
            st.plotly_chart(render_daily_progress_chart(g_daily), use_container_width=True, key="chart_daily_progress")
        else:
            st.caption("Belum ada data harian untuk periode ini.")

        cc1, cc2 = st.columns(2)
        with cc1:
            st.markdown("###### Kontribusi Marketing Corporate vs Retail")
            if not mc_summary.empty:
                fig_mc = render_mc_split_donut(mc_summary)
                if fig_mc:
                    st.plotly_chart(fig_mc, use_container_width=True, key="chart_mc_split_ringkasan")
            else:
                st.caption("Belum ada data.")
        with cc2:
            st.markdown("###### Kontribusi 6 Pilar")
            if not pilar_summary.empty:
                fig_pilar = render_contribution_pie(
                    [_pilar_label(p) for p in pilar_summary["Pilar"]],
                    pilar_summary["Omset"],
                    [PILAR_COLORS.get(p, "#9ca3af") for p in pilar_summary["Pilar"]],
                )
                st.plotly_chart(fig_pilar, use_container_width=True, key="chart_pilar_ringkasan")
            else:
                st.caption("Belum ada data.")

        st.markdown("###### Riwayat 30 Hari Terakhir")
        g_hist = build_daily_history(df_main, selected_branches)
        if not g_hist.empty:
            st.plotly_chart(render_daily_history_chart(g_hist), use_container_width=True, key="chart_daily_history")

        all_insights = generate_all_sales_insights(df_main) if not df_main.empty else []
        if all_insights:
            st.markdown("###### 💡 Insight & Rekomendasi")
            for ins in all_insights[:5]:
                render_structured_insight_card(ins)

with tab2:
    st.subheader(f"🏆 Scoreboard — {periode_label}")
    st.caption(f"Kuartal berjalan: {quarter_period_label} • Total hari: {q_total_hari} • Hari berjalan: {q_hari_berjalan} • Sisa hari: {q_sisa_hari}")

    exp_c1, exp_c2 = st.columns(2)
    with exp_c1:
        pdf_bytes = generate_scoreboard_pdf(sb_dict)
        st.download_button("📄 Export Scoreboard (PDF)", data=pdf_bytes, file_name=f"scoreboard_mflash_{tanggal_acuan}.pdf",
                            mime="application/pdf", key="dl_scoreboard_pdf")
    with exp_c2:
        st.caption("Export JPEG per kategori tersedia di bawah masing-masing tabel.")

    for kat in SCOREBOARD_KATEGORI:
        st.markdown(render_scoreboard_html(sb_dict[kat], kat), unsafe_allow_html=True)
        jpg_bytes = generate_scoreboard_table_image(sb_dict[kat], kat)
        st.download_button(f"🖼️ Export {kat} (JPEG)", data=jpg_bytes,
                            file_name=f"scoreboard_{sanitize_filename(kat)}_{tanggal_acuan}.jpg",
                            mime="image/jpeg", key=f"dl_scoreboard_jpg_{kat}")
        st.markdown("<br/>", unsafe_allow_html=True)

with tab3:
    st.subheader(f"📢 Iklan (Meta Ads) — {periode_label}")
    if df_ads.empty:
        st.info("Belum ada data Iklan. Unggah file export Meta Ads di sidebar.")
    else:
        ads_agg = aggregate_ads_by_branch(df_ads)
        if selected_branches:
            ads_agg = ads_agg[ads_agg["Cabang"].isin(selected_branches)]
        total_spend = ads_agg["AmountSpent"].sum() if "AmountSpent" in ads_agg.columns else 0
        total_results = ads_agg["Results"].sum() if "Results" in ads_agg.columns else 0
        total_reach = ads_agg["Reach"].sum() if "Reach" in ads_agg.columns else 0

        k1, k2, k3 = st.columns(3)
        with k1:
            st.markdown(render_kpi_card("Total Amount Spent", format_rupiah(total_spend), "#7c3aed", "💸"), unsafe_allow_html=True)
        with k2:
            st.markdown(render_kpi_card("Total Results", format_number(total_results), "#2563eb", "🎯"), unsafe_allow_html=True)
        with k3:
            st.markdown(render_kpi_card("Total Reach", format_number(total_reach), "#059669", "📡"), unsafe_allow_html=True)

        if not ads_agg.empty and "AmountSpent" in ads_agg.columns:
            fig_ads = go.Figure()
            fig_ads.add_trace(go.Bar(x=ads_agg["Cabang"], y=ads_agg["AmountSpent"], marker_color="#7c3aed",
                                      text=[format_rupiah(v) for v in ads_agg["AmountSpent"]], textposition="outside"))
            fig_ads.update_layout(height=340, margin=dict(t=20, b=10, l=10, r=10), xaxis_title="Cabang", yaxis_title="Amount Spent")
            st.plotly_chart(fig_ads, use_container_width=True, key="chart_ads_spend")

        st.dataframe(ads_agg, use_container_width=True, hide_index=True)

        ads_insights = generate_ads_insights(ads_agg)
        if ads_insights:
            st.markdown("###### 💡 Insight & Rekomendasi Iklan")
            for ins in ads_insights[:5]:
                render_insight_card(ins)

with tab4:
    st.subheader(f"🚶 Walk-in per Cabang — {quarter_period_label}")
    if df_walkin.empty:
        st.info("Belum ada data Walk-in. Unggah file export Rincian Pengiriman Pesanan di sidebar.")
    else:
        d_walkin = _walkin_ordered(walkin_agg_current)
        st.markdown(render_walkin_table_html(d_walkin, quarter_period_label), unsafe_allow_html=True)

        we1, we2 = st.columns(2)
        with we1:
            jpg_bytes = generate_walkin_table_image(d_walkin, quarter_period_label)
            st.download_button("🖼️ Export Tabel Walk-in (JPEG)", data=jpg_bytes,
                                file_name=f"walkin_{tanggal_acuan}.jpg", mime="image/jpeg", key="dl_walkin_jpg")
        with we2:
            pdf_bytes = generate_walkin_table_pdf(d_walkin, quarter_period_label)
            st.download_button("📄 Export Tabel Walk-in (PDF)", data=pdf_bytes,
                                file_name=f"walkin_{tanggal_acuan}.pdf", mime="application/pdf", key="dl_walkin_pdf")

        if not d_walkin.empty:
            fig_walkin = go.Figure()
            fig_walkin.add_trace(go.Bar(x=d_walkin["Cabang"], y=d_walkin["TotalWalkin"], marker_color="#0f766e",
                                         text=[format_number(v) for v in d_walkin["TotalWalkin"]], textposition="outside"))
            fig_walkin.update_layout(height=340, margin=dict(t=20, b=10, l=10, r=10), xaxis_title="Cabang", yaxis_title="Total Walk-in")
            st.plotly_chart(fig_walkin, use_container_width=True, key="chart_walkin_total")

        wk_marketing_insights = generate_walkin_marketing_insights(d_walkin)
        wk_trend_insights = generate_walkin_insights(walkin_agg_monthly)
        all_wk_insights = wk_marketing_insights + wk_trend_insights
        if all_wk_insights:
            st.markdown("###### 💡 Insight & Rekomendasi Walk-in")
            for ins in all_wk_insights[:5]:
                render_structured_insight_card(ins)

with tab5:
    st.subheader(f"🧩 6 Pilar MFlash — {quarter_period_label}")
    if pilar_summary.empty:
        st.info("Belum ada data 6 Pilar untuk periode ini.")
    else:
        cols_pilar = st.columns(len(pilar_summary))
        for i, (_, r) in enumerate(pilar_summary.iterrows()):
            with cols_pilar[i]:
                st.markdown(render_pilar_kpi_card(r["Pilar"], r["Omset"], r["Qty"]), unsafe_allow_html=True)

        st.markdown("<br/>", unsafe_allow_html=True)
        cpi1, cpi2 = st.columns([1, 1])
        with cpi1:
            fig_pilar2 = render_contribution_pie(
                [_pilar_label(p) for p in pilar_summary["Pilar"]],
                pilar_summary["Omset"],
                [PILAR_COLORS.get(p, "#9ca3af") for p in pilar_summary["Pilar"]],
                title="Kontribusi Omset per Pilar",
            )
            st.plotly_chart(fig_pilar2, use_container_width=True, key="chart_pilar_tab5")
        with cpi2:
            st.markdown("###### Ringkasan per Pilar")
            st.markdown(render_pilar_table_html(pilar_summary), unsafe_allow_html=True)

        st.markdown("###### Detail per Cabang")
        st.markdown(render_pilar_summary_table_html(pilar_by_branch), unsafe_allow_html=True)

        pilar_insights = generate_pilar_insights(pilar_summary)
        if pilar_insights:
            st.markdown("###### 💡 Insight & Rekomendasi")
            for ins in pilar_insights[:5]:
                render_structured_insight_card(ins)

with tab6:
    st.subheader(f"🤝 Kontribusi Marketing Corporate vs Sales Retail — {quarter_period_label}")
    if mc_summary.empty:
        st.info("Belum ada data untuk periode ini.")
    else:
        total_mc = mc_summary["Omset"].sum()
        cols_mc = st.columns(len(mc_summary))
        for i, (_, r) in enumerate(mc_summary.iterrows()):
            with cols_mc[i]:
                st.markdown(render_mc_contribution_card(r["Kelompok"], r["Omset"], total_mc), unsafe_allow_html=True)

        st.markdown("<br/>", unsafe_allow_html=True)
        fig_mc2 = render_mc_split_donut(mc_summary)
        if fig_mc2:
            st.plotly_chart(fig_mc2, use_container_width=True, key="chart_mc_split_tab6")

        st.markdown("###### Detail Marketing Corporate per Sales")
        st.markdown(render_mc_person_table_html(mc_person_table), unsafe_allow_html=True)

        st.markdown("###### Omset Retail per Cabang")
        st.markdown(render_retail_by_branch_table_html(retail_by_branch), unsafe_allow_html=True)
        if not retail_by_branch.empty:
            fig_retail = go.Figure()
            fig_retail.add_trace(go.Bar(x=retail_by_branch["Cabang"], y=retail_by_branch["Omset"], marker_color="#0f766e",
                                         text=[format_rupiah(v) for v in retail_by_branch["Omset"]], textposition="outside"))
            fig_retail.update_layout(height=340, margin=dict(t=20, b=10, l=10, r=10), xaxis_title="Cabang", yaxis_title="Omset Retail")
            st.plotly_chart(fig_retail, use_container_width=True, key="chart_retail_branch")

        mc_insights = generate_mc_insights(mc_summary)
        retail_insights = generate_retail_branch_insights(retail_by_branch)
        combined_insights = mc_insights + retail_insights
        if combined_insights:
            st.markdown("###### 💡 Insight & Rekomendasi")
            for ins in combined_insights[:5]:
                render_structured_insight_card(ins)

st.markdown("---")
st.subheader("📦 Export Laporan Lengkap")
exp1, exp2 = st.columns(2)
with exp1:
    if st.button("📊 Buat Laporan PPTX", key="btn_gen_pptx"):
        pptx_bytes = generate_pptx_report(df_main, sb_dict, walkin_agg_current, pilar_summary, mc_summary, quarter_period_label)
        st.session_state["_pptx_report"] = pptx_bytes
    if st.session_state.get("_pptx_report"):
        st.download_button("⬇️ Download Laporan PPTX", data=st.session_state["_pptx_report"],
                            file_name=f"laporan_mflash_{tanggal_acuan}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="dl_pptx_report")
with exp2:
    if st.button("📄 Buat Laporan PDF", key="btn_gen_pdf"):
        pdf_bytes = generate_pdf_report(df_main, sb_dict, walkin_agg_current, pilar_summary, mc_summary, quarter_period_label)
        st.session_state["_pdf_report"] = pdf_bytes
    if st.session_state.get("_pdf_report"):
        st.download_button("⬇️ Download Laporan PDF", data=st.session_state["_pdf_report"],
                            file_name=f"laporan_mflash_{tanggal_acuan}.pdf", mime="application/pdf",
                            key="dl_pdf_report")
